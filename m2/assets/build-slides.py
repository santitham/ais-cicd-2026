#!/usr/bin/env python3
"""Build Module 2 slides.pptx, matching the Module 1 deck exactly."""

import copy
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette ---
ORANGE   = C(0xFC, 0x78, 0x00)
INK      = C(0x22, 0x22, 0x24)
BODY     = C(0x3F, 0x3F, 0x46)
MUTED    = C(0x77, 0x77, 0x80)
WHITE    = C(0xFF, 0xFF, 0xFF)
HDR_DARK = C(0x40, 0x40, 0x41)
TERM_BG  = C(0x2E, 0x2E, 0x30)
TERM_TX  = C(0xEC, 0xEC, 0xEE)
GRN_HDR  = C(0x2E, 0x9E, 0x5B)
OUT_BG   = C(0x1C, 0x2C, 0x22)
OUT_TX   = C(0xCC, 0xEE, 0xD8)
OK_BG    = C(0xEC, 0xF7, 0xF0)
NOTE_BG  = C(0xFF, 0xFB, 0xAA)
NOTE_LN  = C(0xF5, 0xD3, 0x28)
RED      = C(0xC0, 0x39, 0x2B)
RED_BG   = C(0xFC, 0xEF, 0xEE)
GREY_BG  = C(0xF4, 0xF4, 0xF6)
GREY_LN  = C(0xD8, 0xD8, 0xDE)

L, W = 2.35, 10.40
FOOT = "Day 1 · Foundations of Git, YAML and Databricks  ·  Module 2"

prs = Presentation("bx-slide-template.pptx")

# remove the template's example slides
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)

BLANK = prs.slide_masters[0].slide_layouts[6]
TOTAL = 0  # filled in after the build


# ---------------------------------------------------------------- helpers ---
def new():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h, text, size=12, bold=False, color=BODY, font="Arial",
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, line=None):
    tb = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(s, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, In(l), In(t), In(w), In(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def chrome(s, n):
    box(s, L, 6.99, 8.20, 0.30, FOOT, 7.5, False, MUTED)
    box(s, 11.05, 6.99, 1.70, 0.30, f"{n} / %%TOTAL%%", 7.5, False, MUTED,
        align=PP_ALIGN.RIGHT)


def head(s, title, sub=None, size=25):
    box(s, L, 0.36, W, 0.62, title, size, True, INK)
    if sub:
        box(s, L, 1.02, W, 0.40, sub, 11.5, False, MUTED)


# --------------------------------------------------------------- archetypes -
def s_title(n, kicker, title, sub, meta):
    s = new()
    box(s, L, 2.32, W, 0.30, kicker, 11, True, ORANGE)
    box(s, L, 2.72, W, 1.10, title, 33, True, INK)
    box(s, L, 3.98, W, 0.80, sub, 14, False, BODY, line=1.25)
    rect(s, L, 5.10, 1.30, 0.045, ORANGE)
    box(s, L, 5.44, W, 0.30, meta, 11, False, MUTED)
    box(s, L, 6.62, W, 0.28, "DataOps Mastery — Advanced CI/CD and Pipeline Automation",
        9.5, False, MUTED)
    chrome(s, n)


def s_part(n, kicker, title, sub):
    s = new()
    rect(s, L, 3.02, W, 0.04, ORANGE)
    box(s, L, 2.56, W, 0.34, kicker, 11, True, ORANGE)
    box(s, L, 3.24, W, 1.30, title, 33, True, INK)
    box(s, L, 4.58, W, 0.80, sub, 13, False, MUTED, line=1.25)
    chrome(s, n)


def s_bullets(n, title, sub, items, note=None, size=13.5, gap=11):
    """items: list of (bold_lead_or_None, text)"""
    s = new()
    head(s, title, sub)
    t = 1.68
    for lead, text in items:
        rect(s, L, t + 0.09, 0.075, 0.075, ORANGE, shape=MSO_SHAPE.OVAL)
        if lead:
            box(s, L + 0.28, t - 0.03, W - 0.28, 0.28, lead, size, True, INK)
            tb = box(s, L + 0.28, t + 0.25, W - 0.28, 0.60, text, size - 1.5,
                     False, BODY, line=1.22)
            t += 0.30 + 0.24 * (1 + text.count("\n")) + gap / 100.0 + 0.16
        else:
            box(s, L + 0.28, t - 0.03, W - 0.28, 0.60, text, size, False, BODY,
                line=1.22)
            t += 0.24 * (1 + text.count("\n")) + gap / 100.0 + 0.10
    if note:
        y = max(t + 0.14, 5.95)
        rect(s, L, y, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_runsee(n, title, sub, run, see, run_label="YOU RUN", see_label="YOU SEE",
             see_ok=True, csize=9.5):
    s = new()
    head(s, title, sub)
    hdr2 = GRN_HDR if see_ok else C(0x9E, 0x3A, 0x2E)
    bg2 = OUT_BG if see_ok else C(0x2E, 0x1C, 0x1C)
    tx2 = OUT_TX if see_ok else C(0xF2, 0xD2, 0xCE)
    rect(s, L, 1.62, 5.03, 0.34, HDR_DARK)
    box(s, L + 0.14, 1.62, 4.75, 0.34, run_label, 9, True, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, L, 1.96, 5.03, 4.92, TERM_BG)
    box(s, L + 0.16, 2.08, 4.71, 4.68, run, csize, False, TERM_TX, "Consolas",
        line=1.24)
    rect(s, 7.72, 1.62, 5.03, 0.34, hdr2)
    box(s, 7.86, 1.62, 4.75, 0.34, see_label, 9, True, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 7.72, 1.96, 5.03, 4.92, bg2)
    box(s, 7.88, 2.08, 4.71, 4.68, see, csize, False, tx2, "Consolas", line=1.24)
    chrome(s, n)


def s_code(n, title, sub, code, note=None, csize=11, ch=4.60):
    s = new()
    head(s, title, sub)
    rect(s, L, 1.62, W, ch, TERM_BG)
    box(s, L + 0.16, 1.74, W - 0.32, ch - 0.24, code, csize, False, TERM_TX,
        "Consolas", line=1.26)
    if note:
        y = 1.62 + ch + 0.16
        rect(s, L, y, W, 0.56, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.56, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    chrome(s, n)


NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # "No Style, No Grid"


def _cell_text(cell, text, size, bold, color, font, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.16
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def s_table(n, title, sub, headers, rows, widths, note=None, size=11.5,
            rowh=0.42, mono_cols=()):
    """A real PowerPoint table: editable cells, resizable columns."""
    s = new()
    head(s, title, sub)

    hdr_h = 0.40
    row_hs = [rowh * (1 + max(c.count("\n") for c in row)) for row in rows]
    top = 1.66

    gf = s.shapes.add_table(len(rows) + 1, len(headers),
                            In(L), In(top), In(W), In(hdr_h + sum(row_hs)))
    tbl = gf.table

    # Suppress the built-in blue table style so our own fills show through.
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for sid in tblPr.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"):
        sid.text = NO_STYLE

    for ci, w in enumerate(widths):
        tbl.columns[ci].width = In(w)
    tbl.rows[0].height = In(hdr_h)
    for ri, h in enumerate(row_hs, start=1):
        tbl.rows[ri].height = In(h)

    def style_cell(cell, fill):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = In(0.14)
        cell.margin_right = In(0.10)
        cell.margin_top = In(0.04)
        cell.margin_bottom = In(0.04)

    for ci, htxt in enumerate(headers):
        c = tbl.cell(0, ci)
        style_cell(c, HDR_DARK)
        _cell_text(c, htxt, 10.5, True, WHITE, "Arial")

    for ri, row in enumerate(rows):
        for ci, cell_txt in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            style_cell(c, GREY_BG if ri % 2 == 0 else WHITE)
            col = INK if ci == 0 else BODY
            bold = ci == 0
            if cell_txt.startswith("*"):
                cell_txt, col, bold = cell_txt[1:], ORANGE, True
            _cell_text(c, cell_txt, size, bold, col,
                       "Consolas" if ci in mono_cols else "Arial")

    if note:
        y = min(top + hdr_h + sum(row_hs) + 0.20, 6.14)
        rect(s, L, y, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_challenge(n, num, title, minutes, steps, success):
    s = new()
    rect(s, L, 0.30, W, 1.24, HDR_DARK)
    rect(s, L, 1.54, W, 0.05, ORANGE)
    box(s, L + 0.24, 0.44, 4.00, 0.30, "LAB", 11, True, ORANGE)
    box(s, L + 0.24, 0.76, 7.80, 0.66, f"Challenge {num} — {title}", 21, True, WHITE)
    rect(s, 10.85, 0.46, 1.60, 0.38, ORANGE)
    box(s, 10.85, 0.46, 1.60, 0.38, minutes, 9.5, True, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t = 1.96
    step = min(0.62, (4.90 - 0.10) / max(len(steps), 1))
    for i, txt in enumerate(steps, 1):
        rect(s, L + 0.10, t, 0.34, 0.34, ORANGE, shape=MSO_SHAPE.OVAL)
        box(s, L + 0.10, t, 0.34, 0.34, str(i), 9, True, WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        box(s, L + 0.58, t - 0.04, 9.72, step - 0.06, txt, 12, False, BODY,
            line=1.20)
        t += step
    y = max(t + 0.12, 6.10)
    rect(s, L, y, W, 0.50, OK_BG, GRN_HDR)
    box(s, L + 0.18, y, W - 0.36, 0.50, "Success: " + success, 11.5, False, INK,
        anchor=MSO_ANCHOR.MIDDLE)
    chrome(s, n)


def s_shot(n, title, sub, ref, commands, what):
    """Screenshot placeholder. `ref` matches assets/capture-screenshots.md."""
    s = new()
    head(s, title, sub)
    rect(s, L, 1.62, W, 0.86, TERM_BG)
    box(s, L + 0.16, 1.72, W - 0.32, 0.68, commands, 10.5, False, TERM_TX,
        "Consolas", line=1.24)
    ph = rect(s, L, 2.64, W, 3.42, C(0xFA, 0xFA, 0xFB), GREY_LN)
    box(s, L, 3.66, W, 0.34, f"SCREENSHOT {ref}", 15, True, MUTED,
        align=PP_ALIGN.CENTER)
    box(s, L, 4.06, W, 0.34, "capture on the Ubuntu server — assets/capture-screenshots.md",
        10.5, False, MUTED, align=PP_ALIGN.CENTER)
    rect(s, L, 6.20, W, 0.56, NOTE_BG, NOTE_LN)
    box(s, L + 0.18, 6.20, W - 0.36, 0.56, what, 11.5, False, INK,
        anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    chrome(s, n)


def s_split(n, title, sub, left_title, left_body, right_title, right_body,
            note=None, mono_left=False, mono_right=False):
    s = new()
    head(s, title, sub)
    for x, ttl, bdy, mono in ((L, left_title, left_body, mono_left),
                              (7.72, right_title, right_body, mono_right)):
        rect(s, x, 1.62, 5.03, 0.40, HDR_DARK)
        box(s, x + 0.16, 1.62, 4.71, 0.40, ttl, 10.5, True, WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x, 2.02, 5.03, 3.90, GREY_BG, GREY_LN)
        box(s, x + 0.18, 2.18, 4.67, 3.58, bdy,
            10 if mono else 12.5, False, BODY,
            "Consolas" if mono else "Arial", line=1.28)
    if note:
        rect(s, L, 6.14, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, 6.14, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_end(n):
    s = new()
    box(s, L, 2.70, W, 0.30, "END OF MODULE 2", 11, True, ORANGE)
    box(s, L, 3.10, W, 0.90, "Next: Bundle anatomy and your first deployment",
        30, True, INK)
    box(s, L, 4.14, 8.60, 0.90,
        "Module 3 takes the configuration you validated this afternoon and gives it "
        "a schema, a target, and a deployment command.", 13.5, False, BODY, line=1.28)
    rect(s, L, 5.30, 1.30, 0.045, ORANGE)
    box(s, L, 5.62, W, 0.30, "santitham.pro@kmutt.ac.th", 12, False, MUTED)
    chrome(s, n)


# ================================================================== CONTENT ==
i = 0
def nx():
    global i
    i += 1
    return i

# --- front matter ------------------------------------------------------------
s_title(nx(), "MODULE 2", "YAML and the Databricks workspace",
        "The configuration language every tool this week reads, and the command line "
        "you will drive the workspace from.",
        "Day 1 · afternoon · 3 hours")

s_bullets(nx(), "Module overview", "By the end of this module, you can", [
    (None, "Read any YAML file as nested mappings, sequences and scalars."),
    (None, "Write correctly indented YAML from a set of requirements."),
    (None, "Control how a value is typed, and know when quoting is necessary."),
    (None, "Use multi-line strings, comments, anchors and aliases appropriately."),
    (None, "Distinguish YAML syntax from a tool's schema, and check each separately."),
    (None, "Configure and verify a Databricks CLI profile, and run a job from the terminal."),
], size=14, gap=20)

s_table(nx(), "How this module runs", "Three hours: two of instruction, one at your own keyboard",
        ["", "Topic", "Teach", "You work"],
        [["Part A · 1", "The mental model: trees, mappings, sequences, scalars", "20 min", "6 min"],
         ["Part A · 2", "Writing YAML: one configuration, built up", "22 min", "7 min"],
         ["Part A · 3", "Types, quoting and multi-line strings", "22 min", "10 min"],
         ["Part A · 4", "Correctness, tools, and the mistakes to know", "18 min", "12 min"],
         ["Break", "", "15 min", ""],
         ["Part B", "The workspace and the Databricks CLI", "23 min", "25 min"],
         ["*Total", "*3 hours", "*120 min", "*60 min"]],
        [1.90, 4.70, 1.90, 1.90], rowh=0.36, size=11,
        note="Labs are timed but not gated. If you finish early, the Going further step on each "
             "challenge slide has more.")

s_bullets(nx(), "Why YAML", "You will meet it in every remaining module of this course", [
    ("It describes structure, not behaviour",
     "YAML holds configuration. It contains no logic and executes nothing. The\n"
     "application that reads the file decides what the keys mean."),
    ("It is designed to be read and edited by people",
     "No braces, no commas, no closing tags. That is the whole reason the industry\n"
     "settled on it for configuration rather than JSON or XML."),
    ("It lives in Git, so configuration gets reviewed",
     "A change to a pipeline becomes a diff someone can read, comment on, and revert."),
    ("The name, and the extension",
     "YAML Ain't Markup Language. .yaml is the recommended extension; .yml means the\n"
     "same thing and is more common in practice."),
], size=13, gap=8)

s_table(nx(), "Where this sits",
        "Module 1 ended with a repository. Module 3 needs a configuration you can trust.",
        ["Module", "What it adds", "Written in"],
        [["1 · Git", "The repository, and Databricks Repos", "—"],
         ["*2 · YAML", "*Hand-written configuration, and the checks that verify it", "*YAML"],
         ["3 · Bundles", "databricks.yml, targets, variables", "YAML"],
         ["5 · Testing", "pytest.ini, a test suite", "YAML + Python"],
         ["6 · CI", "Workflow files running on a server", "YAML"],
         ["8 · Production", "Job specifications, smoke tests", "YAML"]],
        [2.10, 5.80, 2.50], rowh=0.44)

# --- Part A1: the mental model ----------------------------------------------
s_part(nx(), "PART A", "YAML", "Three constructs, one tree, and how to write it correctly.")

s_split(nx(), "A YAML document is a tree",
        "The file is an encoding. The tree is the object. Every tool reads the tree.",
        "WHAT YOU WRITE",
        "pipeline:\n"
        "  name: customer_features\n"
        "  enabled: true\n"
        "  owners:\n"
        "    - data-science\n"
        "    - data-engineering\n"
        "  settings:\n"
        "    retries: 3",
        "WHAT THE TOOL RECEIVES",
        "{\n"
        "  \"pipeline\": {\n"
        "    \"name\": \"customer_features\",\n"
        "    \"enabled\": true,\n"
        "    \"owners\": [\"data-science\",\n"
        "                \"data-engineering\"],\n"
        "    \"settings\": {\"retries\": 3}\n"
        "  }\n"
        "}",
        mono_left=True, mono_right=True,
        note="Read the left, picture the right. That habit is most of what this module teaches.")

s_table(nx(), "You already know this data model",
        "YAML has no structure that Python and JSON do not already have",
        ["YAML term", "Python", "JSON", "In the example opposite"],
        [["Mapping", "dict", "object", "pipeline, settings"],
         ["Sequence", "list", "array", "owners"],
         ["Scalar", "str, int, float, bool, None", "string, number, true/false, null",
          "customer_features, true, 3"]],
        [2.30, 2.90, 3.10, 2.10], rowh=0.50, size=11,
        note="If you can read a nested dictionary in a notebook, you can read any YAML file. "
             "The only new thing is the notation.")

s_split(nx(), "The same object, in two notations",
        "If you can write the dictionary, you can write the YAML",
        "PYTHON",
        "config = {\n"
        "    \"name\": \"churn-model\",\n"
        "    \"workers\": 4,\n"
        "    \"enabled\": True,\n"
        "    \"tags\": [\"ml\", \"analytics\"],\n"
        "    \"compute\": {\n"
        "        \"runtime\": \"15.4\"\n"
        "    }\n"
        "}",
        "YAML",
        "name: churn-model\n"
        "workers: 4\n"
        "enabled: true\n"
        "tags:\n"
        "  - ml\n"
        "  - analytics\n"
        "compute:\n"
        "  runtime: \"15.4\"",
        mono_left=True, mono_right=True,
        note="Braces become indentation. Commas disappear. Quotes become optional. Nothing "
             "else changes — including which values need to stay strings.")

s_bullets(nx(), "Scalars", "A single value at a leaf of the tree", [
    ("key: value", "One key, one scalar. This is the entire syntax at a leaf."),
    ("Four kinds", "Strings, numbers, booleans and null. You write characters; the parser\n"
                   "decides which of the four you meant."),
    ("Empty is not one thing", "tags: []  is an empty list.   parameters: {}  is an empty map.\n"
                               "description: null  and  blank:  are both null.   empty: \"\"  is a string."),
], size=13, gap=10,
   note="Four different kinds of nothing. A schema expecting a list will reject three of them.")

s_split(nx(), "Mappings", "A value may itself be a set of key-value pairs",
        "NESTING BY INDENTATION",
        "compute:\n"
        "  runtime: \"15.4\"\n"
        "  workers: 2\n"
        "  autoscale:\n"
        "    min_workers: 2\n"
        "    max_workers: 8",
        "THE RULES",
        "Spaces only. A tab is an error.\n\n"
        "Sibling keys share a column. runtime,\n"
        "workers and autoscale are siblings\n"
        "because all three start at column 3.\n\n"
        "Two spaces per level is the convention.\n"
        "Keys are case-sensitive.",
        mono_left=True,
        note="Indentation is not formatting. It is the only thing that expresses which key "
             "belongs to which parent.")

s_split(nx(), "Sequences", "An ordered list of values",
        "A LIST OF SCALARS",
        "libraries:\n"
        "  - pandas\n"
        "  - scikit-learn\n"
        "  - mlflow",
        "WHAT IT BECOMES",
        "{\n"
        "  \"libraries\": [\n"
        "      \"pandas\",\n"
        "      \"scikit-learn\",\n"
        "      \"mlflow\"\n"
        "  ]\n"
        "}",
        mono_left=True, mono_right=True,
        note="Every `- ` marker in one list must start in the same column. Misalign one and "
             "you get a different structure, usually without an error.")

s_split(nx(), "Sequences of mappings", "The single most common shape in this course",
        "EACH ELEMENT IS A MAP",
        "tasks:\n"
        "  - name: prepare_data\n"
        "    notebook: notebooks/prepare\n"
        "    retries: 2\n"
        "  - name: train_model\n"
        "    notebook: notebooks/train\n"
        "    retries: 1",
        "HOW TO READ IT",
        "`- ` opens a new element.\n\n"
        "notebook and retries belong to the\n"
        "element above them because they are\n"
        "indented to the same column as the text\n"
        "after `- `.\n\n"
        "Two elements here, each with three keys.",
        mono_left=True,
        note="tasks in a bundle, steps in a workflow, and every pipeline definition you write "
             "this week has exactly this shape.")

s_split(nx(), "Reading a file you did not write", "A method that works on any of them",
        "THE FILE",
        "resources:\n"
        "  jobs:\n"
        "    nightly:\n"
        "      name: nightly-scoring\n"
        "      tasks:\n"
        "        - task_key: score\n"
        "          libraries:\n"
        "            - pypi:\n"
        "                package: mlflow",
        "THE METHOD",
        "1  Find the leftmost column. Those keys are\n"
        "   the top level: here, resources alone.\n\n"
        "2  Follow one branch down, ignoring the rest.\n\n"
        "3  At each line ask only: mapping, sequence,\n"
        "   or scalar?\n\n"
        "4  A `- ` means you have entered a list. Count\n"
        "   how many you passed to know how deep you are.\n\n"
        "5  Confirm with  yq -o json  before you edit.",
        mono_left=True,
        note="Four levels of nesting and two sequences. Nobody reads this correctly at a glance; "
             "everybody can read it correctly with a method.")

s_challenge(nx(), 1, "Read a document you did not write", "6 min", [
    "Open fixtures/pipeline.yml. Do not run anything yet.",
    "On paper, draw the tree. Mark every node as a mapping, a sequence, or a scalar.",
    "Write down how many elements the tasks sequence has, and how many keys each element has.",
    "Now check: yq -o json . fixtures/pipeline.yml",
    "Going further: write the equivalent Python literal for the whole document.",
], "your drawing matches the JSON, and you can point to the line that makes owners a list.")

# --- Part A2: writing it -----------------------------------------------------
s_part(nx(), "PART A · 2", "Writing YAML",
       "Building one configuration, one construct at a time.")

s_split(nx(), "Step 1 — A flat mapping", "Start with the values that have no structure",
        "WHAT YOU WRITE",
        "name: churn-model\n"
        "owner: data-science\n"
        "version: \"1.0\"\n"
        "enabled: true\n"
        "retries: 3",
        "NOTES",
        "One key per line, a colon, a space,\n"
        "then the value. The space is required.\n\n"
        "version is quoted because 1.0 would\n"
        "otherwise become the number 1.\n\n"
        "enabled is a real boolean. retries is a\n"
        "real integer. Neither is quoted.",
        mono_left=True)

s_split(nx(), "Step 2 — Group related keys", "Nest them under a parent",
        "WHAT YOU WRITE",
        "project:\n"
        "  name: churn-model\n"
        "  owner: data-science\n"
        "  version: \"1.0\"\n"
        "\n"
        "notifications:\n"
        "  enabled: true\n"
        "  channel: \"#data-alerts\"",
        "NOTES",
        "project: has no value on its own line.\n"
        "Its value is everything indented beneath it.\n\n"
        "The blank line is decoration. It changes\n"
        "nothing about the structure.\n\n"
        "\"#data-alerts\" must be quoted, or the #\n"
        "starts a comment and the value is null.",
        mono_left=True,
        note="A key with an indented block beneath it has that block as its value. A key with "
             "nothing after it has null.")

s_split(nx(), "Step 3 — Add a list", "For values that are ordered and repeatable",
        "WHAT YOU WRITE",
        "tags:\n"
        "  - machine-learning\n"
        "  - customer-analytics\n"
        "\n"
        "environments:\n"
        "  development:\n"
        "    workers: 1\n"
        "  production:\n"
        "    workers: 4",
        "NOTES",
        "tags is a sequence: order matters, names\n"
        "do not.\n\n"
        "environments is a mapping: the names\n"
        "development and production matter,\n"
        "and order does not.\n\n"
        "Choosing between these two is a design\n"
        "decision, not a syntax one.",
        mono_left=True)

s_split(nx(), "Step 4 — A list of mappings", "When each item needs several fields",
        "WHAT YOU WRITE",
        "tasks:\n"
        "  - name: prepare_features\n"
        "    entrypoint: notebooks/prepare\n"
        "    retries: 2\n"
        "\n"
        "  - name: train_model\n"
        "    entrypoint: notebooks/train\n"
        "    retries: 1",
        "NOTES",
        "The `- ` and the first key share a line.\n"
        "Everything after it aligns to the column\n"
        "where that first key begins.\n\n"
        "Get that column wrong by one space and\n"
        "the parser builds a different tree —\n"
        "sometimes without complaining.",
        mono_left=True,
        note="If you write only one shape correctly from memory, make it this one.")

s_split(nx(), "Mapping or sequence?", "A design decision, and the one beginners get wrong",
        "USE A MAPPING WHEN",
        "Each item has a name that means something,\n"
        "and you will refer to it by that name.\n\n"
        "environments:\n"
        "  development:\n"
        "    workers: 1\n"
        "  production:\n"
        "    workers: 4\n\n"
        "Order is irrelevant. Names are the point.",
        "USE A SEQUENCE WHEN",
        "Items are alike, order may matter, and\n"
        "nothing refers to one by name.\n\n"
        "tasks:\n"
        "  - name: prepare\n"
        "    retries: 2\n"
        "  - name: train\n"
        "    retries: 1\n\n"
        "You can add a third without renaming anything.",
        note="A schema usually fixes this choice for you. When it does not, ask whether anything "
             "will ever need to address one item by name.")

s_bullets(nx(), "Comments", "The one part of the file the parser throws away", [
    ("Everything after # on a line is discarded",
     "It never reaches the tool. It exists only for the next person to open the file,\n"
     "which is usually you, some months from now."),
    ("Record the why, not the what",
     "retries: 3  # transient S3 timeouts during the 02:00 peak.\n"
     "The key already says what it is. Only you know why the value is 3."),
    ("A # inside a value needs quoting",
     "channel: \"#data-alerts\" is a string. channel: #data-alerts is null followed by\n"
     "a comment, and nothing will warn you."),
], size=13, gap=10)

s_bullets(nx(), "Conventions worth adopting", "None of these are enforced. All of them help.", [
    ("Two spaces per level, never tabs",
     "Configure your editor once: expand tabs, two-space indent, show whitespace."),
    ("Comment the why, not the what",
     "retries: 3  # transient S3 failures during peak hours. The key already says\n"
     "what it is; only you know why it is 3."),
    ("Quote anything that is not a quantity or a true boolean",
     "Versions, identifiers, dates, country codes, channel names, anything with a colon."),
    ("Keep one document per file, and name it after what it configures",
     "`---` can separate several documents in one file. It rarely helps."),
], size=13, gap=6)

s_bullets(nx(), "Changing a file you did not write", "The procedure that keeps you out of trouble", [
    ("Parse it before you touch it", "yq . on the original. If it already fails, that is not your fault\n"
                                     "and you need to know before you add a second problem."),
    ("Change one value, never the indentation", "Editing a value is almost always safe. Re-indenting a block is\n"
                                                "where structure gets silently rewritten."),
    ("Copy the shape of the line above", "To add a list element, duplicate an existing one and edit the copy.\n"
                                         "You inherit its indentation, which is the part that is easy to get wrong."),
    ("Parse it again, and read the diff", "yq . to confirm it still parses, then read the diff before you save.\n"
                                          "Whitespace-only changes in a diff are worth a second look."),
], size=13, gap=6)

s_bullets(nx(), "Set your editor up once", "Five minutes now, or an afternoon later", [
    ("Expand tabs to two spaces", "In VS Code: \"editor.insertSpaces\": true, \"editor.tabSize\": 2.\n"
                                  "This alone removes an entire class of error permanently."),
    ("Render whitespace", "\"editor.renderWhitespace\": \"all\". Tabs and trailing spaces become visible\n"
                          "instead of being things you have to remember to suspect."),
    ("Install the YAML extension", "The Red Hat YAML extension validates as you type and, for files it\n"
                                   "recognises, checks them against a published schema."),
    ("Which gives you layers 1 and 2 in the editor", "You still run the command-line checks. But most mistakes never survive\n"
                                                     "long enough to reach them."),
], size=13, gap=6)

# --- Part A3: types and quoting ---------------------------------------------
s_bullets(nx(), "Types are resolved, not declared",
          "The author writes characters; the parser decides what they mean", [
    ("There is no type annotation in YAML",
     "You cannot write that version is a string. Nothing in the syntax lets you say so."),
    ("The parser pattern-matches the token",
     "An unquoted value is compared against a table of patterns. The first match wins,\n"
     "and the decision is not reported anywhere."),
    ("Quoting is how you take the decision back",
     "\"1.0\" is a string because you said so. 1.0 is a float because the parser said so."),
], note="This is the one idea in Part A worth memorising. Most surprises in configuration "
        "files come from it.")

s_table(nx(), "Implicit type resolution", "What an unquoted value becomes",
        ["Written", "Resolves to", "Type"],
        [["42", "42", "integer"],
         ["1.0", "1.0", "float"],
         ["3.10", "3.1", "float  — the zero is lost"],
         ["true, True, TRUE", "true", "boolean"],
         ["null, ~, or nothing", "null", "null"],
         ["2026-08-16", "a date object", "timestamp"],
         ["*\"3.10\"", "*3.10", "*string"]],
        [3.10, 4.60, 2.70], rowh=0.38,
        note="The last row is the defence, and it is the only one that works in every parser "
             "and every version.")

s_split(nx(), "When to quote", "A short rule that covers almost every case",
        "QUOTE IT",
        "version: \"1.0\"\n"
        "python_version: \"3.10\"\n"
        "date_label: \"2026-08-16\"\n"
        "country: \"NO\"\n"
        "channel: \"#data-alerts\"\n"
        "message: \"key: value\"\n"
        "expression: \"${var.catalog}\"",
        "LEAVE IT UNQUOTED",
        "workers: 4\n"
        "threshold: 0.85\n"
        "enabled: true\n"
        "description: null\n"
        "name: churn-model\n"
        "entrypoint: notebooks/train",
        mono_left=True, mono_right=True,
        note="If the value is a quantity or a true boolean, leave it. Otherwise quote it. "
             "The cost of quoting something unnecessarily is zero.")

s_bullets(nx(), "Expressions are just strings", "To YAML, at least", [
    ("Other systems embed their own syntax",
     "${var.catalog} in a bundle, ${{ github.sha }} in a workflow. YAML has no idea what\n"
     "either of these means."),
    ("YAML's only job is to deliver the string intact",
     "The consuming tool substitutes the value afterwards, long after parsing is over."),
    ("Which is why they need quoting",
     "A value beginning with { is flow-mapping syntax to YAML. Quote it and the problem\n"
     "disappears entirely."),
], note="The semantics of these expressions belong to Modules 3 and 6. Today they are strings "
        "that must survive parsing unchanged.")

s_challenge(nx(), 2, "Predict the type", "7 min", [
    "Write down the resolved type of each: 42 · 1.0 · 3.10 · 0755 · true · null · 1e3 · 2026-08-16 · \"3.10\"",
    "Commit to all nine before running anything. A guess you did not write down does not count.",
    "Check each with:  echo 'k: <value>' > t.yml && yq '.k | type' t.yml",
    "For any you got wrong, find the row of the table that explains it.",
    "Then the question that matters: which of the nine would you have quoted, before today?",
], "you can state the rule deciding each of the nine, not merely the answer.")

# --- Part A4: multi-line -----------------------------------------------------
s_split(nx(), "Multi-line strings", "Two indicators, and the difference matters",
        "LITERAL  |  — line breaks preserved",
        "description: |\n"
        "  Train the model.\n"
        "  Register the best result.\n"
        "  Notify the owning team.\n"
        "\n"
        "→ \"Train the model.\\nRegister the\n"
        "   best result.\\nNotify...\"",
        "FOLDED  >  — line breaks become spaces",
        "description: >\n"
        "  This description is written across\n"
        "  several lines but is read as one\n"
        "  paragraph.\n"
        "\n"
        "→ \"This description is written across\n"
        "   several lines but is read as...\"",
        mono_left=True, mono_right=True,
        note="Both produce an ordinary string. Chomping indicators control the final newline: "
             "| keeps one, |- strips it, |+ keeps all.")

s_bullets(nx(), "Choosing between them", "One rule, and one case where it matters a great deal", [
    ("Use | for anything where lines are meaningful",
     "Shell commands, SQL, formatted text, anything a reader will scan line by line."),
    ("Use > for prose that happens to be wrapped",
     "A long description you broke across lines only to keep the file narrow."),
    ("For a run: block, always use |",
     "A run: value is a script handed to an interpreter as one string. Under > the lines\n"
     "join, and two commands arrive as one malformed command."),
], note="Nothing rejects the wrong choice. The fault appears in a runner log, as a shell error "
        "about an argument you never wrote.")

s_runsee(nx(), "The same lines under each indicator", "What the consumer actually receives",
         "literal: |\n"
         "  line one\n"
         "  line two\n"
         "\n"
         "folded: >\n"
         "  line one\n"
         "  line two\n"
         "\n"
         "literal_strip: |-\n"
         "  line one\n"
         "  line two\n"
         "\n"
         "plain_multiline: line one\n"
         "  line two\n"
         "\n"
         "$ yq -o json . fixtures/scalars.yml",
         "{\n"
         "  \"literal\":\n"
         "      \"line one\\nline two\\n\",\n"
         "\n"
         "  \"folded\":\n"
         "      \"line one line two\\n\",\n"
         "\n"
         "  \"literal_strip\":\n"
         "      \"line one\\nline two\",\n"
         "\n"
         "  \"plain_multiline\":\n"
         "      \"line one line two\"\n"
         "}\n"
         "\n"
         "A plain multi-line scalar folds,\n"
         "exactly like >.")

s_shot(nx(), "A run: block under both indicators",
       "The same two commands, and the two strings the shell would receive", "S18",
       "$ cat fixtures/runblock.yml\n$ yq -o json . fixtures/runblock.yml",
       "Highlight the second run value: two commands folded onto one line, which is why | is "
       "the right choice for scripts.")

# --- Part A5: reuse ----------------------------------------------------------
s_split(nx(), "Anchors and aliases", "Define once, expand many times",
        "SOURCE",
        "default_compute: &default_compute\n"
        "  workers: 2\n"
        "  runtime: \"15.4\"\n"
        "\n"
        "development:\n"
        "  compute: *default_compute\n"
        "production:\n"
        "  compute: *default_compute",
        "AFTER PARSING",
        "development:\n"
        "  compute:\n"
        "    workers: 2\n"
        "    runtime: \"15.4\"\n"
        "production:\n"
        "  compute:\n"
        "    workers: 2\n"
        "    runtime: \"15.4\"",
        mono_left=True, mono_right=True,
        note="Expansion happens in the parser. The consumer receives a fully expanded tree and "
             "cannot tell an anchor was used.")

s_bullets(nx(), "Use them sparingly", "They remove repetition and add indirection", [
    ("They give exact copies, not variants",
     "Merge keys (<<:) would let you inherit a block and override one field. GitHub\n"
     "Actions does not implement them, and support elsewhere varies by parser."),
    ("They do not cross files",
     "An anchor is resolved only within the document that defines it."),
    ("Overuse makes a file harder to read than the repetition did",
     "A reader has to hold the anchor definition in mind while reading every use of it."),
    ("Tool-native reuse is usually better",
     "Targets and variables in bundles (Modules 3 and 4). Reusable workflows and\n"
     "composite actions in GitHub Actions (Module 7)."),
], size=13, gap=6)

# --- Part A6: correctness ----------------------------------------------------
s_bullets(nx(), "Three layers of correctness",
          "A file can pass one and fail the next, and the messages look nothing alike", [
    ("1 · YAML syntax", "Can a parser turn these characters into a tree? If not, nothing else\n"
                        "matters and the parser will tell you the line."),
    ("2 · The tool's schema", "Are the expected keys present, in the expected places, holding the\n"
                              "expected kinds of value? A YAML parser has no opinion about this."),
    ("3 · Runtime behaviour", "Do the referenced notebooks, credentials, clusters and commands\n"
                              "actually exist and work? Only running it answers this."),
], note="favorite_planet: Saturn is perfectly valid YAML. It is not a valid workflow, and it "
        "is not a valid bundle. Layer 1 passed; layer 2 did not.")

s_table(nx(), "Checking each layer", "The instruments, and what each is blind to",
        ["Layer", "Command", "Catches", "Blind to"],
        [["1 · Syntax", "yq .", "Tabs, bad indent, missing space", "Anything that parses"],
         ["1 · Types", "yq '<path> | type'", "3.10 became a float", "Anything you don't ask about"],
         ["2 · Style", "yamllint", "Duplicate keys, truthy values", "Tool-specific structure"],
         ["2 · Schema", "actionlint", "Unknown keys, wrong inputs", "Whether the file exists"],
         ["3 · Runtime", "run it", "Everything else", "Nothing — but it is slow"]],
        [1.90, 2.40, 3.20, 2.90], rowh=0.40, size=10.5,
        note="Work down the list. Each check is cheaper and more precise than the one below it, "
             "and each is blind to what the next one sees.")

s_bullets(nx(), "yq", "The instrument you will use most", [
    ("yq . file.yml", "Parse and print the tree back. A failure means the file is not valid\n"
                      "YAML, and the message carries a line and a column."),
    ("yq -o json . file.yml", "The same tree as JSON. The clearest way to see what a tool receives."),
    ("yq '<path> | type' file.yml", "Report the resolved type of one value: !!str, !!int, !!float, !!bool."),
], note="Do not install yq with apt. Ubuntu's package of that name is a Python wrapper around "
        "jq, reports version 0.0.0, and has incompatible syntax. Install the Go binary.")

s_shot(nx(), "yq on a valid file", "Parsing, and the same document as JSON", "S23",
       "$ yq . fixtures/pipeline.yml\n$ yq -o json . fixtures/pipeline.yml",
       "yq preserves flow style: [main] stays [main]. It is reprinting the tree, not the file.")

# --- Part A7: mistakes, consolidated ----------------------------------------
s_table(nx(), "Six mistakes worth recognising",
        "Ordered by how hard they are to notice, not how often they happen",
        ["Written", "What you get", "Caught by"],
        [["A tab in the indentation", "Nothing — the file does not parse", "yq, immediately"],
         ["name:training-job", "Parse error, or a scalar", "yq, immediately"],
         ["environment: twice", "One value silently discarded", "yamllint"],
         ["A misaligned  -  marker", "One list element instead of two", "nothing, usually"],
         ["python-version: 3.10", "The number 3.1", "nothing — you must look"],
         ["A key indented one level too deep", "A correct value in the wrong place", "the tool's schema"]],
        [3.90, 3.40, 3.10], rowh=0.40, size=10.5,
        note="The first two stop you. The last four do not, which is why the type check and the "
             "schema check are worth running rather than trusting your reading.")

s_split(nx(), "The two that parse cleanly", "Same file, two parsers, three different answers",
        "WHAT IS IN THE FILE",
        "on:\n"
        "  push:\n"
        "    branches: [main]\n"
        "python-version: 3.10\n"
        "enabled: no\n"
        "country: NO",
        "WHAT EACH PARSER RETURNS",
        "                  yq v4      PyYAML\n"
        "                 (1.2)       (1.1)\n"
        "on:             \"on\"        true\n"
        "python-version   3.1         3.1\n"
        "enabled         \"no\"        false\n"
        "country         \"NO\"        false",
        mono_left=True, mono_right=True,
        note="YAML 1.1 read yes/no/on/off as booleans; 1.2 removed that. Which specification "
             "applies depends on the parser reading your file, not on the file.")

s_shot(nx(), "One file, two parsers", "The demonstration, run on your own machine", "S33",
       "$ yq -o json . fixtures/two-parsers.yml\n"
       "$ python3 -c \"import yaml,json;print(json.dumps(yaml.safe_load(open('fixtures/two-parsers.yml')),indent=2))\"",
       "The defence is the same one from the quoting rule: quote the value and both parsers "
       "agree.")

s_bullets(nx(), "Reading a parser error", "The message is more useful than it first appears", [
    ("found character that cannot start any token", "A tab. Almost always a tab."),
    ("mapping values are not allowed in this context",
     "A colon where a plain value was expected — usually a missing space after a colon,\n"
     "or a line indented one level too deep beneath a scalar."),
    ("did not find expected key",
     "Inconsistent indentation. The parser reached a column it cannot place."),
    ("The reported line is where it noticed, not where you erred",
     "The fault is on that line or the one above it. Read both before editing anything."),
], size=13, gap=6)

s_challenge(nx(), 3, "Break it, then read the message", "10 min", [
    "Copy fixtures/pipeline.yml to work.yml. Restore it to this state before each step.",
    "Replace the leading spaces on one line with a tab. Run yq and read the message.",
    "Remove the space after a colon. Predict the message before you run it.",
    "Indent one `- ` marker by one extra space. Run yq — then explain why there is no error.",
    "Duplicate a key with a different value. Find the one tool that objects.",
], "for each, you can state what the parser said and what it actually meant.")

s_challenge(nx(), 4, "Author a configuration from requirements", "12 min", [
    "Start from assets/skeleton.yml — five empty keys and nothing else.",
    "Add a project name and owner; development and production environments with different worker counts.",
    "Add two tasks, each with a name, entrypoint and retry count; a boolean controlling notifications.",
    "Add a multi-line description with line breaks preserved; two tags; a version whose exact text must survive.",
    "Validate with yq, then swap with your neighbour and have them state your structure aloud.",
], "yq -o json shows the tree you intended, and your neighbour read it the same way you did.")

# --- Part A8: bridge ---------------------------------------------------------
s_split(nx(), "The same structures, in the files you meet next",
        "Read these as trees. The schemas are Modules 3 and 6; the shapes are today.",
        "DATABRICKS ASSET BUNDLE",
        "bundle:\n"
        "  name: customer-churn\n"
        "resources:\n"
        "  jobs:\n"
        "    training_job:\n"
        "      name: training-job\n"
        "      tasks:\n"
        "        - task_key: train\n"
        "          notebook_task:\n"
        "            notebook_path: ./src/train",
        "GITHUB ACTIONS WORKFLOW",
        "name: Validate configuration\n"
        "on:\n"
        "  pull_request:\n"
        "jobs:\n"
        "  validate:\n"
        "    runs-on: ubuntu-latest\n"
        "    steps:\n"
        "      - name: Run validation\n"
        "        run: example-command validate",
        mono_left=True, mono_right=True,
        note="Same three constructs in both. In each: jobs is a mapping, the named job is a "
             "mapping, tasks and steps are sequences, and every element of them is a mapping.")

s_bullets(nx(), "Exit ticket", "Answer these before the break", [
    ("What are the three constructs a YAML document is built from?",
     "And what is each one called in Python?"),
    ("Why can a file be valid YAML and still be rejected by a bundle or a workflow?",
     "Name the layer at which each of those two checks operates."),
    ("Which value on this slide needs quoting, and why?    version: 1.0    workers: 4    country: NO",
     "State the rule, not just the answer."),
    ("After an indentation error, what do you inspect first?",
     "And why is the line the parser reports often not the line you need to edit?"),
], size=13, gap=8)

# --- Part B ------------------------------------------------------------------
s_part(nx(), "PART B", "The workspace and the CLI",
       "The same three layers, applied to a second tool.")

s_bullets(nx(), "The workspace", "What is in it, and which parts this course touches", [
    ("Workspace files", "Notebooks, Python files and folders, including the Git folders you\n"
                        "created in Module 1."),
    ("Compute", "Clusters and warehouses. Created, started, stopped and billed\n"
                "independently of the code that runs on them."),
    ("Jobs", "A named unit of work: one or more tasks, each with something to run\n"
             "and compute to run it on. Module 3 generates these from YAML."),
    ("Unity Catalog", "Catalogs, schemas, tables. You need write access to one catalog;\n"
                      "the rest is outside this course."),
], size=13, gap=8)

s_table(nx(), "Compute types", "Three kinds, and they are not interchangeable",
        ["", "All-purpose", "Jobs compute", "Serverless"],
        [["Created by", "A person, once", "The job, per run", "The platform"],
         ["Lifetime", "Until stopped", "The run only", "The run only"],
         ["Configuration", "Edited in the UI", "Declared in a file", "Mostly fixed"],
         ["Start-up", "Already warm", "Several minutes", "Seconds"],
         ["Cost profile", "Highest", "Lower", "Per query"]],
        [2.30, 2.70, 2.90, 2.50], rowh=0.44, size=11)

s_bullets(nx(), "Which compute for automation",
          "Why a jobs cluster, and not the one you already have running", [
    ("An all-purpose cluster is shared mutable state",
     "Someone installs a library, changes a Spark setting, or restarts it on a different\n"
     "runtime. None of that is recorded in your repository."),
    ("Which makes runs irreproducible",
     "A pipeline that passes today can fail tomorrow for a reason that appears in no\n"
     "commit and no diff."),
    ("Jobs compute is declared, not maintained",
     "The specification lives in a file, under review, in Git. That is the property\n"
     "automation needs, and it is the reason to prefer it."),
    ("Serverless trades control for latency",
     "No start-up wait, less control over the runtime. Reasonable for short tasks."),
], size=13, gap=6)

s_bullets(nx(), "The CLI and the workspace", "A local program issuing remote calls", [
    ("The same relation Git has to GitHub", "From Module 1: some operations are local, some cross the network.\n"
                                            "Confusing the two causes most beginner errors."),
    ("Except that almost everything here is remote", "The CLI holds no local state beyond a configuration file. Every\n"
                                                     "command this afternoon is an API call."),
    ("Two programs share the name", "The modern CLI is one Go binary, v0.205 or later. The legacy\n"
                                    "databricks-cli is a Python package with incompatible subcommands."),
], note="`databricks --version` distinguishes them: `Databricks CLI v0.2xx.x` against "
        "`databricks-cli, version 0.18.0`.")

s_table(nx(), "The command surface", "Six groups cover everything this course does",
        ["Group", "What it operates on", "Used in"],
        [["databricks auth / configure", "Profiles and credentials", "Today"],
         ["databricks workspace", "Files and folders in the workspace", "Today"],
         ["databricks jobs", "Job definitions and runs", "Today, and Module 8"],
         ["databricks bundle", "Deploy a whole project from a file", "Modules 3, 4 and 7"],
         ["databricks clusters", "Compute lifecycle", "Rarely, once bundles exist"],
         ["databricks fs", "DBFS paths", "Occasionally"]],
        [3.60, 4.10, 2.70], rowh=0.40, size=11,
        note="Every group takes --profile and --output json. Learn those two flags and the rest "
             "of the CLI is discoverable with --help.")

s_shot(nx(), "Version check", "Confirming which of the two programs you have", "S53",
       "$ databricks --version",
       "If this prints `databricks-cli, version 0.18.x`, remove it before continuing. Nothing "
       "in this course will work with the legacy package.")

s_bullets(nx(), "Authentication", "A host, and a credential", [
    ("Personal access token", "Generated from User Settings → Developer → Access tokens.\n"
                              "Set a short expiry: a token with no expiry outlives the course."),
    ("It is a bearer credential", "Anyone holding it acts as you, with your permissions. There is no\n"
                                  "second factor and no confirmation step."),
    ("Where it must never go", "Not committed. Not pasted into a notebook. Not written into a\n"
                               "workflow file. Not sent in a chat message."),
    ("This is temporary", "Module 7 replaces it with a service principal, which is what a\n"
                          "pipeline should authenticate as."),
], size=13, gap=8)

s_split(nx(), "The configuration file", "One file, several named profiles",
        "~/.databrickscfg",
        "[dev]\n"
        "host  = https://adb-1234.5.azuredatabricks.net\n"
        "token = dapi****************************\n"
        "\n"
        "[staging]\n"
        "host  = https://adb-6789.0.azuredatabricks.net\n"
        "token = dapi****************************",
        "SELECTING ONE",
        "databricks <cmd> --profile dev\n"
        "\n"
        "export DATABRICKS_CONFIG_PROFILE=dev\n"
        "\n"
        "The dev / staging / prod split you create\n"
        "here becomes bundle targets in Module 4,\n"
        "and deployment environments in Module 7.",
        mono_left=True,
        note="The file lives outside any repository by design. That, not .gitignore, is what "
             "keeps the token out of your history.")

s_shot(nx(), "configure, and the file it writes", "Two commands and their result", "S56",
       "$ databricks configure --host https://<your-workspace> --profile dev\n$ cat ~/.databrickscfg",
       "Capture with a throwaway token and revoke it immediately afterwards. A redaction box "
       "invites the question of what was underneath it.")

s_bullets(nx(), "Verify before use", "databricks current-user me", [
    ("This is the schema check for the CLI", "It confirms three things at once: the configuration parsed, the host\n"
                                             "is reachable, and the credential was accepted."),
    ("Run it before any real work", "It converts a class of confusing downstream failures into one clear\n"
                                    "failure, at the moment you know what you just changed."),
    ("Read the failures deliberately", "A wrong host, an expired token and a missing profile each produce a\n"
                                       "different message. Recognising them is the skill."),
], note="Same discipline as running yq before you rely on a file: check the cheap thing first, "
        "where the answer is unambiguous.")

s_challenge(nx(), 5, "Configure and verify", "8 min", [
    "Generate a personal access token with a short expiry. Note the expiry date now.",
    "Run: databricks configure --host <your-workspace> --profile dev",
    "cat ~/.databrickscfg and confirm the profile name and host are what you expect.",
    "Verify: databricks current-user me --profile dev | jq .userName",
    "Now break it deliberately: change the host to a wrong value, re-run, and read the error.",
], "a working dev profile, and you recognise the error a wrong host produces.")

s_runsee(nx(), "Reading the output with jq", "Every command can return JSON; jq extracts one field",
         "# human-readable by default\n"
         "$ databricks workspace list \"$WS\" \\\n"
         "    --profile dev\n"
         "\n"
         "\n"
         "# machine-readable on request\n"
         "$ databricks workspace list \"$WS\" \\\n"
         "    --profile dev --output json\n"
         "\n"
         "\n"
         "\n"
         "\n"
         "# one field out of the structure\n"
         "$ databricks workspace list \"$WS\" \\\n"
         "    --profile dev --output json \\\n"
         "  | jq -r '.[].path'",
         "hello_world\n"
         "\n"
         "\n"
         "\n"
         "\n"
         "[\n"
         "  {\n"
         "    \"object_type\": \"NOTEBOOK\",\n"
         "    \"path\": \"/Workspace/.../hello_world\",\n"
         "    \"language\": \"PYTHON\"\n"
         "  }\n"
         "]\n"
         "\n"
         "\n"
         "/Workspace/.../hello_world",
         csize=9.5,
         see_label="YOU SEE")

s_bullets(nx(), "Three jq patterns cover this course", "You do not need to learn jq properly", [
    ("jq .", "Pretty-print the whole structure. Run this first, always, to see what\n"
             "you are working with before writing a path."),
    ("jq -r '.field.subfield'", "Extract one value. -r strips the surrounding quotes, which matters\n"
                                "when you assign the result to a shell variable."),
    ("jq -r '.[].name'", "Extract one field from every element of a list."),
], note="jq -r '.state.result_state' is the only one you will actually need to remember. It "
        "appears again in Module 8.")

s_bullets(nx(), "The job lifecycle", "Five operations, entirely from the terminal", [
    ("workspace import", "Upload a file into the workspace at a path you choose."),
    ("Create the job", "In the UI this afternoon. Note the job identifier from the URL."),
    ("jobs run-now", "Trigger a run. Returns a run identifier immediately; it does not wait."),
    ("jobs get-run", "Poll. Returns life_cycle_state while running, result_state at the end."),
    ("jobs get-run-output", "Retrieve what the task printed."),
], note="Module 3 replaces this whole sequence, including the UI step, with one command: "
        "databricks bundle deploy.")

s_runsee(nx(), "Polling and the state object", "What a run reports while it runs, and after",
         "$ RUN_ID=$(databricks jobs run-now \\\n"
         "    --job-id $JOB_ID --profile dev \\\n"
         "    | jq -r .run_id)\n"
         "\n"
         "$ databricks jobs get-run \\\n"
         "    --run-id \"$RUN_ID\" --profile dev \\\n"
         "    | jq .state\n"
         "\n"
         "\n"
         "\n"
         "# sixty seconds later\n"
         "\n"
         "\n"
         "\n"
         "\n"
         "# the single field CI actually needs\n"
         "$ databricks jobs get-run \\\n"
         "    --run-id \"$RUN_ID\" --profile dev \\\n"
         "    | jq -r '.state.result_state'",
         "\n\n\n\n"
         "{\n"
         "  \"life_cycle_state\": \"RUNNING\",\n"
         "  \"state_message\": \"In run\"\n"
         "}\n"
         "\n"
         "\n"
         "\n"
         "{\n"
         "  \"life_cycle_state\": \"TERMINATED\",\n"
         "  \"result_state\": \"SUCCESS\",\n"
         "  \"state_message\": \"\"\n"
         "}\n"
         "\n"
         "\n"
         "SUCCESS")

s_shot(nx(), "A complete job lifecycle", "Import through result, in one terminal", "S59",
       "$ databricks workspace mkdirs \"$WS\" --profile dev\n"
       "$ databricks workspace import --file hello_world.py --format SOURCE --language PYTHON --profile dev \"$WS/hello_world\"\n"
       "$ RUN_ID=$(databricks jobs run-now --job-id $JOB_ID --profile dev | jq -r .run_id)\n"
       "$ databricks jobs get-run --run-id \"$RUN_ID\" --profile dev | jq -r '.state.result_state'",
       "That last command is the smoke test in Module 8. Name it as such here, so it is "
       "recognised when it reappears on Thursday.")

s_challenge(nx(), 6, "Full lifecycle without the UI", "17 min", [
    "Create the workspace folder and import assets/hello_world.py into it from the CLI.",
    "Confirm it arrived: databricks workspace list \"$WS\" --profile dev --output json | jq",
    "Create the job in the UI — the one step still done by hand — and capture its job id.",
    "Trigger it with run-now, capture the run id with jq, and poll until it terminates.",
    "Retrieve the output. Then answer: which of these five steps survives Module 3?",
], "a SUCCESS result state obtained by command, and jq used to extract at least three fields.")

s_table(nx(), "Common errors", "Symptom, cause, recovery",
        ["Symptom", "Cause", "Fix"],
        [["yq: unknown command", "apt's yq — a jq wrapper", "Install the Go binary"],
         ["mapping values are not allowed", "Missing space after a colon", "Read that line and the one above"],
         ["cannot configure default credentials", "No profile selected", "Add --profile, or export DATABRICKS_CONFIG_PROFILE"],
         ["403 Forbidden", "Token expired or insufficient", "Re-generate the token; check workspace access"],
         ["jq: error — not defined", "Querying a path that is absent", "Pipe through jq . first and read the structure"]],
        [3.60, 3.10, 3.70], rowh=0.44, size=10.5)

s_bullets(nx(), "Before Module 3", "The state you should be in tomorrow morning", [
    (None, "You can read a nested YAML document and state its structure without running it."),
    (None, "You quote versions, identifiers, dates and country codes by reflex."),
    (None, "You know which of | and > to use for a multi-line command, and why."),
    (None, "You can name the three layers of correctness and the tool that checks each."),
    (None, "yq is installed, and it is the Go binary, not the apt package."),
    (None, "databricks current-user me returns your username against a dev profile."),
    (None, "You have run a job and read its result state without opening the workspace UI."),
], size=13.5, gap=16)

s_end(nx())

TOTAL = i

# stamp the real total into every page number
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and "%%TOTAL%%" in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = r.text.replace("%%TOTAL%%", str(TOTAL))

prs.save("slides.pptx")
print(f"wrote slides.pptx — {TOTAL} slides")
