#!/usr/bin/env python3
"""Build Module 3 slides.pptx, matching the Module 1 and Module 2 decks."""

import copy
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette ---
ORANGE   = C(0xFC, 0x78, 0x00)
INK      = C(0x22, 0x22, 0x24)
BODY     = C(0x3F, 0x3F, 0x46)
MUTED    = C(0x77, 0x77, 0x80)
WHITE    = C(0xFF, 0xFF, 0xFF)
HDR_DARK = C(0x40, 0x40, 0x41)
TERM_BG  = C(0x2E, 0x2E, 0x30)
TERM_TX  = C(0xEC, 0xEC, 0xEE)
GRN_HDR  = C(0x2E, 0x9E, 0x5B)
OUT_BG   = C(0x1C, 0x2C, 0x22)
OUT_TX   = C(0xCC, 0xEE, 0xD8)
OK_BG    = C(0xEC, 0xF7, 0xF0)
NOTE_BG  = C(0xFF, 0xFB, 0xAA)
NOTE_LN  = C(0xF5, 0xD3, 0x28)
RED      = C(0xC0, 0x39, 0x2B)
RED_BG   = C(0xFC, 0xEF, 0xEE)
GREY_BG  = C(0xF4, 0xF4, 0xF6)
GREY_LN  = C(0xD8, 0xD8, 0xDE)

L, W = 2.35, 10.40
FOOT = "Day 2 · Declarative Automation Bundles  ·  Module 3"

prs = Presentation("bx-slide-template.pptx")

# remove the template's example slides
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)

BLANK = prs.slide_masters[0].slide_layouts[6]
TOTAL = 0  # filled in after the build


# ---------------------------------------------------------------- helpers ---
def new():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h, text, size=12, bold=False, color=BODY, font="Arial",
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, line=None):
    tb = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(s, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, In(l), In(t), In(w), In(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def chrome(s, n):
    box(s, L, 6.99, 8.20, 0.30, FOOT, 7.5, False, MUTED)
    box(s, 11.05, 6.99, 1.70, 0.30, f"{n} / %%TOTAL%%", 7.5, False, MUTED,
        align=PP_ALIGN.RIGHT)


def head(s, title, sub=None, size=25):
    box(s, L, 0.36, W, 0.62, title, size, True, INK)
    if sub:
        box(s, L, 1.02, W, 0.40, sub, 11.5, False, MUTED)


# --------------------------------------------------------------- archetypes -
def s_title(n, kicker, title, sub, meta):
    s = new()
    box(s, L, 2.32, W, 0.30, kicker, 11, True, ORANGE)
    box(s, L, 2.72, W, 1.10, title, 33, True, INK)
    box(s, L, 3.98, W, 0.80, sub, 14, False, BODY, line=1.25)
    rect(s, L, 5.10, 1.30, 0.045, ORANGE)
    box(s, L, 5.44, W, 0.30, meta, 11, False, MUTED)
    box(s, L, 6.62, W, 0.28, "DataOps Mastery — Advanced CI/CD and Pipeline Automation",
        9.5, False, MUTED)
    chrome(s, n)


def s_part(n, kicker, title, sub):
    s = new()
    rect(s, L, 3.02, W, 0.04, ORANGE)
    box(s, L, 2.56, W, 0.34, kicker, 11, True, ORANGE)
    box(s, L, 3.24, W, 1.30, title, 33, True, INK)
    box(s, L, 4.58, W, 0.80, sub, 13, False, MUTED, line=1.25)
    chrome(s, n)


def s_bullets(n, title, sub, items, note=None, size=13.5, gap=11):
    """items: list of (bold_lead_or_None, text)"""
    s = new()
    head(s, title, sub)
    t = 1.68
    for lead, text in items:
        rect(s, L, t + 0.09, 0.075, 0.075, ORANGE, shape=MSO_SHAPE.OVAL)
        if lead:
            box(s, L + 0.28, t - 0.03, W - 0.28, 0.28, lead, size, True, INK)
            tb = box(s, L + 0.28, t + 0.25, W - 0.28, 0.60, text, size - 1.5,
                     False, BODY, line=1.22)
            t += 0.30 + 0.24 * (1 + text.count("\n")) + gap / 100.0 + 0.16
        else:
            box(s, L + 0.28, t - 0.03, W - 0.28, 0.60, text, size, False, BODY,
                line=1.22)
            t += 0.24 * (1 + text.count("\n")) + gap / 100.0 + 0.10
    if note:
        y = max(t + 0.14, 5.95)
        rect(s, L, y, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_runsee(n, title, sub, run, see, run_label="YOU RUN", see_label="YOU SEE",
             see_ok=True, csize=9.5):
    s = new()
    head(s, title, sub)
    hdr2 = GRN_HDR if see_ok else C(0x9E, 0x3A, 0x2E)
    bg2 = OUT_BG if see_ok else C(0x2E, 0x1C, 0x1C)
    tx2 = OUT_TX if see_ok else C(0xF2, 0xD2, 0xCE)
    rect(s, L, 1.62, 5.03, 0.34, HDR_DARK)
    box(s, L + 0.14, 1.62, 4.75, 0.34, run_label, 9, True, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, L, 1.96, 5.03, 4.92, TERM_BG)
    box(s, L + 0.16, 2.08, 4.71, 4.68, run, csize, False, TERM_TX, "Consolas",
        line=1.24)
    rect(s, 7.72, 1.62, 5.03, 0.34, hdr2)
    box(s, 7.86, 1.62, 4.75, 0.34, see_label, 9, True, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 7.72, 1.96, 5.03, 4.92, bg2)
    box(s, 7.88, 2.08, 4.71, 4.68, see, csize, False, tx2, "Consolas", line=1.24)
    chrome(s, n)


def s_code(n, title, sub, code, note=None, csize=11, ch=4.60):
    s = new()
    head(s, title, sub)
    rect(s, L, 1.62, W, ch, TERM_BG)
    box(s, L + 0.16, 1.74, W - 0.32, ch - 0.24, code, csize, False, TERM_TX,
        "Consolas", line=1.26)
    if note:
        y = 1.62 + ch + 0.16
        rect(s, L, y, W, 0.56, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.56, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    chrome(s, n)


NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # "No Style, No Grid"


def _cell_text(cell, text, size, bold, color, font, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.16
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def s_table(n, title, sub, headers, rows, widths, note=None, size=11.5,
            rowh=0.42, mono_cols=()):
    """A real PowerPoint table: editable cells, resizable columns."""
    s = new()
    head(s, title, sub)

    hdr_h = 0.40
    row_hs = [rowh * (1 + max(c.count("\n") for c in row)) for row in rows]
    top = 1.66

    gf = s.shapes.add_table(len(rows) + 1, len(headers),
                            In(L), In(top), In(W), In(hdr_h + sum(row_hs)))
    tbl = gf.table

    # Suppress the built-in blue table style so our own fills show through.
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for sid in tblPr.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"):
        sid.text = NO_STYLE

    for ci, w in enumerate(widths):
        tbl.columns[ci].width = In(w)
    tbl.rows[0].height = In(hdr_h)
    for ri, h in enumerate(row_hs, start=1):
        tbl.rows[ri].height = In(h)

    def style_cell(cell, fill):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = In(0.14)
        cell.margin_right = In(0.10)
        cell.margin_top = In(0.04)
        cell.margin_bottom = In(0.04)

    for ci, htxt in enumerate(headers):
        c = tbl.cell(0, ci)
        style_cell(c, HDR_DARK)
        _cell_text(c, htxt, 10.5, True, WHITE, "Arial")

    for ri, row in enumerate(rows):
        for ci, cell_txt in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            style_cell(c, GREY_BG if ri % 2 == 0 else WHITE)
            col = INK if ci == 0 else BODY
            bold = ci == 0
            if cell_txt.startswith("*"):
                cell_txt, col, bold = cell_txt[1:], ORANGE, True
            _cell_text(c, cell_txt, size, bold, col,
                       "Consolas" if ci in mono_cols else "Arial")

    if note:
        y = min(top + hdr_h + sum(row_hs) + 0.20, 6.14)
        rect(s, L, y, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_challenge(n, num, title, minutes, steps, success):
    s = new()
    rect(s, L, 0.30, W, 1.24, HDR_DARK)
    rect(s, L, 1.54, W, 0.05, ORANGE)
    box(s, L + 0.24, 0.44, 4.00, 0.30, "LAB", 11, True, ORANGE)
    box(s, L + 0.24, 0.76, 7.80, 0.66, f"Challenge {num} — {title}", 21, True, WHITE)
    rect(s, 10.85, 0.46, 1.60, 0.38, ORANGE)
    box(s, 10.85, 0.46, 1.60, 0.38, minutes, 9.5, True, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t = 1.96
    step = min(0.62, (4.90 - 0.10) / max(len(steps), 1))
    for i, txt in enumerate(steps, 1):
        rect(s, L + 0.10, t, 0.34, 0.34, ORANGE, shape=MSO_SHAPE.OVAL)
        box(s, L + 0.10, t, 0.34, 0.34, str(i), 9, True, WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        box(s, L + 0.58, t - 0.04, 9.72, step - 0.06, txt, 12, False, BODY,
            line=1.20)
        t += step
    y = max(t + 0.12, 6.10)
    rect(s, L, y, W, 0.50, OK_BG, GRN_HDR)
    box(s, L + 0.18, y, W - 0.36, 0.50, "Success: " + success, 11.5, False, INK,
        anchor=MSO_ANCHOR.MIDDLE)
    chrome(s, n)


def s_shot(n, title, sub, ref, commands, what):
    """Screenshot placeholder. `ref` matches assets/capture-screenshots.md."""
    s = new()
    head(s, title, sub)
    rect(s, L, 1.62, W, 0.86, TERM_BG)
    box(s, L + 0.16, 1.72, W - 0.32, 0.68, commands, 10.5, False, TERM_TX,
        "Consolas", line=1.24)
    ph = rect(s, L, 2.64, W, 3.42, C(0xFA, 0xFA, 0xFB), GREY_LN)
    box(s, L, 3.66, W, 0.34, f"SCREENSHOT {ref}", 15, True, MUTED,
        align=PP_ALIGN.CENTER)
    box(s, L, 4.06, W, 0.34, "capture on the Ubuntu server — assets/capture-screenshots.md",
        10.5, False, MUTED, align=PP_ALIGN.CENTER)
    rect(s, L, 6.20, W, 0.56, NOTE_BG, NOTE_LN)
    box(s, L + 0.18, 6.20, W - 0.36, 0.56, what, 11.5, False, INK,
        anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    chrome(s, n)


def s_split(n, title, sub, left_title, left_body, right_title, right_body,
            note=None, mono_left=False, mono_right=False):
    s = new()
    head(s, title, sub)
    for x, ttl, bdy, mono in ((L, left_title, left_body, mono_left),
                              (7.72, right_title, right_body, mono_right)):
        rect(s, x, 1.62, 5.03, 0.40, HDR_DARK)
        box(s, x + 0.16, 1.62, 4.71, 0.40, ttl, 10.5, True, WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x, 2.02, 5.03, 3.90, GREY_BG, GREY_LN)
        box(s, x + 0.18, 2.18, 4.67, 3.58, bdy,
            10 if mono else 12.5, False, BODY,
            "Consolas" if mono else "Arial", line=1.28)
    if note:
        rect(s, L, 6.14, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, 6.14, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


# ================================================================== CONTENT ==
i = 0
def nx():
    global i
    i += 1
    return i


def s_end(n):
    s = new()
    box(s, L, 2.70, W, 0.30, "END OF MODULE 3", 11, True, ORANGE)
    box(s, L, 3.10, W, 0.90, "Next: Build, deploy, run, iterate", 30, True, INK)
    box(s, L, 4.14, 8.60, 0.90,
        "Module 4 writes the resource files you read this morning: jobs with several "
        "tasks, cluster specifications, and one bundle promoted across targets.",
        13.5, False, BODY, line=1.28)
    rect(s, L, 5.30, 1.30, 0.045, ORANGE)
    box(s, L, 5.62, W, 0.30, "santitham.pro@kmutt.ac.th", 12, False, MUTED)
    chrome(s, n)


# --- front matter ------------------------------------------------------------
s_title(nx(), "MODULE 3", "Bundle anatomy and your first deployment",
        "The file that holds the definition of a job, and the commands that put it "
        "into a workspace and take it out again.",
        "Day 2 · morning · 3 hours")

s_bullets(nx(), "Module overview", "By the end of this module, you can", [
    (None, "State what a bundle contains and what a deployment of it writes to the workspace."),
    (None, "Read databricks.yml and name the function of each top-level key."),
    (None, "Predict the configuration the CLI resolves for a target, and confirm it."),
    (None, "Generate a project with bundle init and account for every file it produces."),
    (None, "Validate, plan, deploy, run, destroy and redeploy from the command line."),
    (None, "Identify which copy of a notebook a deployed job is executing."),
], size=14, gap=20)

s_table(nx(), "How this module runs",
        "Three hours: 105 minutes of instruction, 60 at your own keyboard",
        ["", "Topic", "Teach", "You work"],
        [["Part A", "What a bundle is, and what it replaces", "14 min", "—"],
         ["Part B", "The anatomy of databricks.yml, one key at a time", "33 min", "7 min"],
         ["Break", "", "15 min", ""],
         ["Part C · 1", "bundle init and the project it generates", "14 min", "8 min"],
         ["Part C · 2", "Validation and configuration resolution", "8 min", "10 min"],
         ["Part C · 3", "Plan, deploy, and deployment state", "12 min", "13 min"],
         ["Part C · 4", "Running, summarising and destroying", "9 min", "22 min"],
         ["Part D", "Two copies, six failure modes, exit ticket", "8 min", "—"],
         ["*Total", "*3 hours", "*105 min", "*60 min"]],
        [1.90, 4.70, 1.90, 1.90], rowh=0.34, size=11,
        note="Every challenge from Challenge 4 onward depends on the one before it. If a step "
             "does not complete, say so rather than working ahead.")

s_table(nx(), "Where this sits",
        "Module 2 left you driving a job by hand. Module 4 writes the files this module reads.",
        ["Module", "What it adds", "Artefact"],
        [["2 · YAML and the CLI", "A verified profile, and a job created by five commands", "~/.databrickscfg"],
         ["*3 · Bundle anatomy", "*The declaration of that job, and one deployment command", "*databricks.yml"],
         ["4 · Build and deploy", "Multi-task jobs, cluster specs, promotion across targets", "resources/*.yml"],
         ["5 · Testing", "Tests over the code the bundle deploys", "tests/"],
         ["7 · GitHub Actions", "The same commands, run on a server against a service principal", "workflow files"]],
        [3.00, 5.30, 2.10], rowh=0.44, size=11)

# --- Part A ------------------------------------------------------------------
s_part(nx(), "PART A", "What a bundle is",
       "The problem it solves, what it contains, and the commands that act on it.")

s_bullets(nx(), "Terminology and versions",
          "Two names for one thing, and the version this course requires", [
    ("Databricks Asset Bundles were renamed Declarative Automation Bundles in March 2026",
     "The rename is not breaking. The command is still databricks bundle and no\n"
     "configuration file changes. Material written before March 2026, including this\n"
     "course's syllabus, uses the older name."),
    ("The CLI is at version 1.x",
     "This course requires v1.3.0 or later. Verify with databricks --version before\n"
     "the first lab; v0.2xx accepts bundle commands but behaves differently."),
    ("New bundles use the direct deployment engine",
     "Since v1.3.0 the CLI applies changes through the Databricks SDK rather than\n"
     "Terraform. The consequence you will see is the name of the state file."),
], size=13, gap=10)

s_bullets(nx(), "Where workspace state comes from",
          "The job you created in Module 2, considered as an object someone else has to maintain", [
    (None, "The notebook arrived through workspace import, from a path on your laptop."),
    (None, "The job was created by an API call whose body you assembled on the command line."),
    (None, "The compute was chosen when the job was created, by cluster id."),
    (None, "Recreating the job in a second workspace requires all three, and the workspace "
           "records none of them."),
    (None, "Reviewing a change to the job before it happens is not possible, because there "
           "is nothing to read."),
], size=13.5, gap=18,
    note="This is the state Module 2 deliberately left you in. Part A names what is missing.")

s_table(nx(), "What automated deployment requires",
        "Three properties, and the mechanism that supplies each",
        ["Requirement", "What it means", "Supplied by"],
        [["Reproducibility", "The same input produces the same workspace state",
          "A declaration held in a file"],
         ["Environment parity", "Development and production differ only where declared",
          "One file, several targets"],
         ["Reviewability", "A change is legible before it is applied",
          "A text file, and a plan command"]],
        [2.40, 5.10, 2.90], rowh=0.46, size=11.5,
        note="A bundle is the unit that provides all three. The rest of Part A is what it "
             "contains and what it does.")

s_split(nx(), "A bundle is a source tree with a declaration",
        "The directory is the unit that is deployed",
        "ON DISK",
        "my-project/\n"
        "  databricks.yml      the declaration\n"
        "  pyproject.toml\n"
        "  resources/\n"
        "    sample_job.job.yml\n"
        "  src/\n"
        "    sample_notebook.ipynb\n"
        "    my_project/main.py\n"
        "  tests/",
        "WHAT databricks.yml DECLARES",
        "bundle:      the identity of this project\n"
        "include:     which other files to merge\n"
        "variables:   values that differ by target\n"
        "targets:     where it can be deployed\n"
        "resources:   what the workspace should\n"
        "             contain\n"
        "\n"
        "Part B takes these five keys in order.",
        mono_left=True, mono_right=True)

s_code(nx(), "What a deployment consists of",
       "Four things are written under one root path, derived from the bundle name and the target",
       "/Workspace/Users/<you>/.bundle/<bundle name>/<target>/\n"
       "\n"
       "    files/        the synchronised source tree\n"
       "    artifacts/    anything built during deployment, such as a wheel\n"
       "    resources/    the definitions the workspace was given\n"
       "    state/        what the CLI created, recorded so it can compute a difference\n"
       "\n"
       "and, outside the root path, the jobs and pipelines themselves.",
       csize=12, ch=3.90,
       note="The job appears under Workflows like any other. The four directories above are "
            "what makes it reproducible.")

s_table(nx(), "The command surface", "databricks bundle, with the six used today marked",
        ["Command", "What it does", "Used in"],
        [["*init", "*Generate a project from a template", "*today"],
         ["*validate", "*Resolve the configuration and check it against the schema", "*today"],
         ["*plan", "*Report the actions a deployment would take", "*today"],
         ["*deploy", "*Build, upload, create or update resources, record state", "*today"],
         ["*run", "*Start a job or pipeline and report its terminal state", "*today"],
         ["*destroy", "*Delete the resources and files this bundle created", "*today"],
         ["summary, open", "List deployed resources; open one in a browser", "today"],
         ["generate", "Write configuration from a job that already exists", "Module 4"],
         ["sync, schema, deployment", "One-way file sync; JSON schema; state operations", "reference"]],
        [2.70, 5.60, 2.10], rowh=0.34, size=10.5)

s_table(nx(), "The lifecycle, against Module 2's sequence",
        "The left column is what you ran yesterday afternoon",
        ["Module 2", "Module 3"],
        [["workspace import to upload the notebook", "bundle deploy"],
         ["jobs create with an assembled JSON body", "bundle deploy"],
         ["jobs run-now with the job id", "bundle run <key>"],
         ["jobs get-run in a polling loop", "bundle run reports the terminal state"],
         ["jobs get-run-output with the task run id", "bundle run prints the run URL"],
         ["delete by hand, or leave it behind", "bundle destroy"]],
        [5.20, 5.20], rowh=0.42, size=11.5, mono_cols=(0, 1),
        note="Five commands and a bespoke polling loop become three commands that read one file.")

# --- Part B ------------------------------------------------------------------
s_part(nx(), "PART B", "The anatomy of databricks.yml",
       "One key at a time. Every construct here appears in the file the template "
       "generates after the break.")

s_code(nx(), "The smallest bundle that deploys",
       "Three keys. The next fifteen slides add one mechanism each.",
       "bundle:\n"
       "  name: m3demo\n"
       "\n"
       "targets:\n"
       "  dev:\n"
       "    default: true\n"
       "    workspace:\n"
       "      host: https://adb-1234567890.7.azuredatabricks.net\n"
       "\n"
       "resources:\n"
       "  jobs:\n"
       "    sample_job:\n"
       "      name: sample_job\n"
       "      tasks:\n"
       "        - task_key: notebook_task\n"
       "          notebook_task:\n"
       "            notebook_path: ./src/sample_notebook.ipynb",
       csize=11.5, ch=5.00)

s_bullets(nx(), "bundle: the identity of the project",
          "Two keys, both consequential", [
    ("name", "Identifies the bundle, and appears in every path derived from it. Two bundles\n"
             "with the same name, deployed by the same user, occupy the same root path."),
    ("uuid", "Generated once, at initialisation, and written into the file. It identifies the\n"
             "bundle to the workspace independently of its name, so the name can change."),
], size=13.5, gap=14,
    note="Nothing else belongs under bundle:. Deployment settings belong to a target.")

s_code(nx(), "targets: a named destination",
       "A target names a workspace and the settings that apply when deploying there",
       "bundle:\n"
       "  name: m3demo\n"
       "\n"
       "targets:\n"
       "  dev:\n"
       "    default: true\n"
       "    workspace:\n"
       "      host: https://adb-1234567890.7.azuredatabricks.net\n"
       "\n"
       "  prod:\n"
       "    workspace:\n"
       "      host: https://adb-1234567890.7.azuredatabricks.net\n"
       "      root_path: /Workspace/Users/you@example.com/.bundle/${bundle.name}/${bundle.target}",
       csize=11, ch=4.10,
       note="dev and prod are names chosen by convention, not keywords. Everything declared "
            "outside targets: is common to all of them.")

s_bullets(nx(), "Selecting a target", "Which workspace a command acts on", [
    ("databricks bundle deploy -t prod",
     "The -t, or --target, flag selects one target by name. Every bundle command accepts it."),
    ("default: true",
     "Marks the target used when the flag is omitted. The generated template puts it on dev."),
    ("Nothing prevents prod from being the default",
     "Which is the argument for naming the target explicitly in any command whose effect\n"
     "you would not want to discover afterwards."),
], size=13.5, gap=14)

s_table(nx(), "Workspace paths are derived, not written",
        "From the bundle name and the selected target, resolved here for target dev",
        ["Key", "Resolved value", "Holds"],
        [["root_path", ".../Users/<you>/.bundle/m3demo/dev", "everything below"],
         ["file_path", "<root_path>/files", "the synchronised source tree"],
         ["artifact_path", "<root_path>/artifacts", "built artifacts, such as a wheel"],
         ["resource_path", "<root_path>/resources", "generated resource definitions"],
         ["state_path", "<root_path>/state", "the deployment state and metadata"]],
        [1.90, 4.70, 3.80], rowh=0.42, size=11, mono_cols=(0, 1),
        note="Any of the five may be overridden. The generated prod target overrides root_path, "
             "for a reason given on the production mode slide.")

s_bullets(nx(), "resources: what the workspace should contain",
          "A mapping from resource type, to resource key, to the definition", [
    (None, "resources.jobs.<key> declares a job. resources.pipelines.<key> declares a pipeline."),
    (None, "The type list also covers clusters, model serving endpoints, dashboards, schemas, "
           "volumes and apps."),
    (None, "A resource is declared once and deployed to every target, unless a target overrides it."),
    (None, "This module deploys the resources the template writes. Module 4 authors them."),
], size=13.5, gap=18)

s_code(nx(), "A job resource",
       "One task, taken from the generated project and reduced to its required keys",
       "resources:\n"
       "  jobs:\n"
       "    sample_job:                              # the resource key\n"
       "      name: sample_job                       # the label shown in the workspace\n"
       "\n"
       "      tasks:\n"
       "        - task_key: notebook_task\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/sample_notebook.ipynb",
       csize=11.5, ch=3.40,
       note="notebook_path is relative to the file that declares it, so a resource file and the "
            "sources it names move together.")

s_split(nx(), "Resource key and resource name",
        "Two identifiers on the same object, with different consumers",
        "THE KEY",
        "sample_job\n"
        "\n"
        "Used by the CLI and by references from\n"
        "other resources:\n"
        "\n"
        "  databricks bundle run sample_job\n"
        "  ${resources.jobs.sample_job.id}\n"
        "\n"
        "Changing it detaches the resource from\n"
        "its deployment state.",
        "THE NAME",
        "sample_job\n"
        "\n"
        "The label shown under Workflows.\n"
        "\n"
        "Development mode rewrites it to\n"
        "[dev santitham_pro] sample_job\n"
        "and leaves the key untouched.\n"
        "\n"
        "Changing it affects nothing else.",
        mono_left=True, mono_right=True)

s_code(nx(), "include: composing the configuration from several files",
       "Every matching file is merged into one configuration before anything else happens",
       "# databricks.yml\n"
       "bundle:\n"
       "  name: m3demo\n"
       "\n"
       "include:\n"
       "  - resources/*.yml\n"
       "\n"
       "# resources/sample_job.job.yml\n"
       "resources:\n"
       "  jobs:\n"
       "    sample_job:\n"
       "      name: sample_job\n"
       "      ...",
       csize=11.5, ch=4.30,
       note="The generated project puts one resource in each file, so databricks.yml holds only "
            "bundle-level configuration.")

s_code(nx(), "variables: declaring a parameter",
       "Declared once at the top level, assigned per target",
       "variables:\n"
       "  catalog:\n"
       "    description: The catalog to use\n"
       "  schema:\n"
       "    description: The schema to use\n"
       "\n"
       "targets:\n"
       "  dev:\n"
       "    variables:\n"
       "      catalog: training\n"
       "      schema: ${workspace.current_user.short_name}\n"
       "  prod:\n"
       "    variables:\n"
       "      catalog: training\n"
       "      schema: prod",
       csize=11, ch=3.95,
       note="A variable used but not declared is an error. A variable declared with neither a "
            "default nor a target assignment is prompted for.")

s_table(nx(), "Where a variable's value comes from",
        "Four sources, in increasing order of precedence",
        ["Source", "Form", "Typical use"],
        [["The declaration", "default: value under variables:", "a value that rarely changes"],
         ["The selected target", "variables: block inside the target", "the value for this environment"],
         ["The environment", "DATABRICKS_BUNDLE_VAR_catalog=x", "a value a server supplies"],
         ["The command line", "--var catalog=x", "one deployment that differs"]],
        [2.50, 4.40, 3.50], rowh=0.44, size=11, mono_cols=(1,),
        note="Module 4 promotes one bundle across targets using the second row; Module 7 uses the "
             "third from a GitHub Actions runner.")

s_table(nx(), "Substitution namespaces",
        "Four appear in the generated file, and they do not all resolve at the same time",
        ["Reference", "Resolves to", "Resolved during"],
        [["${bundle.name}, ${bundle.target}", "bundle metadata", "configuration resolution"],
         ["${workspace.file_path}", "a derived workspace path", "configuration resolution"],
         ["${workspace.current_user.short_name}", "the authenticated user, normalised",
          "configuration resolution"],
         ["${var.catalog}", "the variable's value for this target", "configuration resolution"],
         ["${resources.pipelines.etl.id}", "the id of a deployed resource", "*deployment"]],
        [4.00, 3.40, 3.00], rowh=0.40, size=10.5, mono_cols=(0,),
        note="The last row is the only one still written as a reference in validate output, "
             "because the resource it names does not exist until it is deployed.")

s_bullets(nx(), "mode: development",
          "A target may declare a mode, which applies computed settings on top of the declaration", [
    (None, "The declaration is unchanged. The mode adds presets, which the CLI computes and "
           "applies during resolution."),
    (None, "The purpose is that a development deployment cannot be mistaken for, or collide "
           "with, a production one."),
    (None, "Presets can be overridden individually, and a setting on a resource outranks a preset."),
], size=13.5, gap=18,
    note="The next slide is the resolved presets, read out of validate -o json rather than "
         "out of the documentation.")

s_table(nx(), "The presets development mode computes",
        "Present in the resolved configuration, absent from every file",
        ["Preset", "Value", "Effect on the deployment"],
        [["name_prefix", "[dev santitham_pro] ", "every resource name is prefixed"],
         ["tags", "dev: santitham_pro", "jobs and pipelines carry the tag"],
         ["trigger_pause_status", "PAUSED", "schedules and triggers do not fire"],
         ["jobs_max_concurrent_runs", "4", "repeated runs during iteration do not queue"],
         ["pipelines_development", "true", "pipelines run in development mode"]],
        [3.30, 2.80, 4.30], rowh=0.40, size=10.5, mono_cols=(0, 1),
        note="The deployment lock is also disabled, and --cluster-id may override the declared "
             "compute. Neither is permitted in production mode.")

s_bullets(nx(), "mode: production",
          "Computes no presets, and enforces constraints instead", [
    (None, "Paths must not be user-specific, which is why the generated prod target writes "
           "root_path explicitly."),
    (None, "run_as and permissions must be stated, so that ownership does not depend on who "
           "ran the command."),
    (None, "Pipelines must not be marked as development."),
    (None, "--cluster-id is refused outright."),
    (None, "The current Git branch is checked against the one the target declares, unless "
           "--force is passed."),
], size=13, gap=14,
    note="Every constraint above turns a class of accidental production change into a refusal "
         "at the point of deployment.")

s_table(nx(), "One file, two targets",
        "The same sample_job resource, resolved for each target",
        ["", "target dev", "target prod"],
        [["job name", "[dev santitham_pro] sample_job", "sample_job"],
         ["trigger", "pause_status: PAUSED", "pause_status: UNPAUSED"],
         ["tags", "dev: santitham_pro", "none"],
         ["max concurrent runs", "4", "declared value only"],
         ["root path", ".../.bundle/m3demo/dev", ".../.bundle/m3demo/prod"]],
        [2.60, 4.10, 3.70], rowh=0.42, size=10.5, mono_cols=(1, 2),
        note="Nothing in the resource file differs between the two columns. The difference is "
             "the target, and the mode it declares.")

s_bullets(nx(), "How the configuration is assembled",
          "Five stages, in this order, every time a bundle command runs", [
    ("1 · Read", "databricks.yml is parsed."),
    ("2 · Merge", "Every file named by include: is merged into the same configuration."),
    ("3 · Overlay", "The selected target's settings are applied over the common configuration."),
    ("4 · Preset", "The presets implied by the target's mode are computed and applied."),
    ("5 · Substitute", "References are replaced by values, except resource ids."),
], size=13, gap=6,
    note="The result is what the workspace receives, and it is not the text of any file in the "
         "project. Module 2 made the same distinction between a file and its parse.")

s_challenge(nx(), 1, "Predict the resolved configuration", "7 min", [
    "A fifteen-line bundle is on the handout. Target dev is selected, and it declares "
    "mode: development.",
    "Without running anything, write down: the resolved job name, the resolved value of "
    "the ${var.catalog} reference, the trigger pause status, and the workspace root path.",
    "For each answer, write the stage of the assembly sequence that produced it.",
    "Keep the paper. You check these answers against the CLI in Challenge 2.",
], "four written answers, each attached to one of the five assembly stages.")

# --- Part C ------------------------------------------------------------------
s_part(nx(), "PART C", "From template to running job",
       "Generate a project, read it, validate it, deploy it, run it, and remove it.")

s_table(nx(), "bundle init and the built-in templates",
        "databricks bundle init <template>, or a directory path, or a Git repository URL",
        ["Template", "Generates", "Requires"],
        [["*default-python", "*A job with notebook, wheel and pipeline tasks",
          "*uv, and a writable catalog"],
         ["default-minimal", "The same databricks.yml skeleton, no resources", "nothing further"],
         ["default-sql", "A job running .sql files on a SQL warehouse", "a SQL warehouse"],
         ["default-scala", "A JAR job", "a JDK and sbt"],
         ["dbt-sql", "A dbt project driven by a job", "dbt"],
         ["pydabs", "default-python, with resources declared in Python", "uv"],
         ["mlops-stacks", "A full MLOps project with several environments", "considerable setup"]],
        [2.60, 5.20, 2.60], rowh=0.36, size=10.5, mono_cols=(0,),
        note="An organisation distributes its own template as a Git repository, which is how a "
             "team standardises the shape of every new project.")

s_shot(nx(), "Initialisation requires an authenticated workspace",
       "bundle init is not an offline command",
       "S34",
       "$ databricks bundle init default-python\n"
       "Error: default auth: cannot configure default credentials",
       "The template calls the metastore assignment endpoint for the default catalog, and the "
       "SCIM endpoint for your user name. Without credentials it stops before writing a file.")

s_shot(nx(), "What default-python generates",
       "One command, thirty files. The tree is the reference for Challenge 2.",
       "S35",
       "$ databricks bundle init default-python\n"
       "$ tree -a -L 2 m3demo",
       "databricks.yml, pyproject.toml, .gitignore, resources/ with a job and a pipeline, src/ "
       "with a package and a notebook, tests/, fixtures/, and editor configuration.")

s_code(nx(), "The generated databricks.yml",
       "Every construct on this slide was introduced in Part B",
       "bundle:\n"
       "  name: m3demo\n"
       "  uuid: e9720c45-37a9-4b3a-af26-bf5a6cd159a3\n"
       "\n"
       "include:\n"
       "  - resources/*.yml\n"
       "\n"
       "artifacts:\n"
       "  python_artifact:\n"
       "    type: whl\n"
       "    build: uv build --wheel\n"
       "\n"
       "variables:\n"
       "  catalog: {description: The catalog to use}\n"
       "  schema:  {description: The schema to use}\n"
       "\n"
       "targets:\n"
       "  dev:  {mode: development, default: true, ...}\n"
       "  prod: {mode: production, ...}",
       csize=10.5, ch=4.45,
       note="artifacts: is the one key Part B did not cover. It declares how a build product is "
            "produced, here a wheel built by uv.")

s_code(nx(), "The generated job resource",
       "resources/sample_job.job.yml — three tasks, read rather than written",
       "resources:\n"
       "  jobs:\n"
       "    sample_job:\n"
       "      name: sample_job\n"
       "      trigger:\n"
       "        periodic: {interval: 1, unit: DAYS}\n"
       "      parameters:\n"
       "        - {name: catalog, default: '${var.catalog}'}\n"
       "        - {name: schema,  default: '${var.schema}'}\n"
       "      tasks:\n"
       "        - task_key: notebook_task\n"
       "          notebook_task: {notebook_path: ../src/sample_notebook.ipynb}\n"
       "        - task_key: python_wheel_task\n"
       "          depends_on: [{task_key: notebook_task}]\n"
       "          python_wheel_task: {package_name: m3demo, entry_point: main}\n"
       "        - task_key: refresh_pipeline\n"
       "          depends_on: [{task_key: notebook_task}]\n"
       "          pipeline_task: {pipeline_id: '${resources.pipelines.m3demo_etl.id}'}",
       csize=10, ch=4.05,
       note="The last task references a resource declared in another file. That reference is the "
            "one that resolves at deployment rather than during resolution.")

s_table(nx(), "The remaining generated files",
        "What each is for, and whether it affects the deployment",
        ["File or directory", "Purpose", "Affects deploy"],
        [["pyproject.toml", "Package metadata, the main entry point, dev dependencies", "yes"],
         [".gitignore", "Excludes .databricks/, build/, dist/ and explorations", "no"],
         ["src/", "The package, the notebook, and the pipeline transformations", "yes"],
         ["tests/", "pytest tests and a conftest.py, used in Module 5", "no"],
         ["fixtures/", "Sample data referenced by the tests", "no"],
         [".vscode/", "Editor settings and recommended extensions", "no"],
         ["AGENTS.md, CLAUDE.md", "Instructions for coding agents", "no"]],
        [2.80, 5.60, 2.00], rowh=0.36, size=10.5, mono_cols=(0,))

s_challenge(nx(), 2, "Initialise and inventory", "8 min", [
    "Run databricks bundle init default-python. Give the project the name your instructor "
    "supplies, and accept the default catalog it proposes.",
    "Produce the file tree and compare it against assets/expected-bundle-tree.txt. Account "
    "for any difference.",
    "Locate in the generated files: the bundle name, the include pattern, one variable "
    "declaration, and the notebook task.",
    "Run databricks bundle validate -o json and check your four Challenge 1 answers.",
], "a tree that matches the reference, and each Challenge 1 answer confirmed or corrected.")

s_bullets(nx(), "bundle validate",
          "Three checks in one command, and the third needs the network", [
    (None, "The assembled configuration matches the bundle schema: keys exist, and values have "
           "the type the schema declares."),
    (None, "Every reference resolves. An undeclared variable or a misspelled namespace fails here."),
    (None, "The workspace named by the target is reachable, and the credential is accepted."),
    (None, "This is level 3 of the four levels of verification from Module 2, applied to a "
           "second tool. Level 4 is still deployment."),
], size=13, gap=14,
    note="Validation says nothing about whether the workspace will accept the operation. That is "
         "the distinction the failure-mode table returns to.")

s_runsee(nx(), "The resolved configuration",
         "bundle validate -o json prints the configuration after all five assembly stages",
         "$ databricks bundle validate -o json \\\n"
         "    | jq '{presets, variables}'\n"
         "\n"
         "\n"
         "# as written in databricks.yml:\n"
         "\n"
         "variables:\n"
         "  catalog:\n"
         "    description: The catalog to use\n"
         "\n"
         "targets:\n"
         "  dev:\n"
         "    mode: development\n"
         "    variables:\n"
         "      catalog: training\n"
         "\n"
         "# presets appears in no file at all",
         "{\n"
         "  \"presets\": {\n"
         "    \"jobs_max_concurrent_runs\": 4,\n"
         "    \"name_prefix\": \"[dev santitham_pro] \",\n"
         "    \"pipelines_development\": true,\n"
         "    \"tags\": { \"dev\": \"santitham_pro\" },\n"
         "    \"trigger_pause_status\": \"PAUSED\"\n"
         "  },\n"
         "  \"variables\": {\n"
         "    \"catalog\": {\n"
         "      \"default\": \"training\",\n"
         "      \"description\": \"The catalog to use\",\n"
         "      \"value\": \"training\"\n"
         "    }\n"
         "  }\n"
         "}",
         csize=9.5)

s_runsee(nx(), "Local paths become workspace paths",
         "Path rewriting happens during resolution, before anything is deployed",
         "$ databricks bundle validate -o json \\\n"
         "    | jq '.resources.jobs.sample_job\n"
         "          .tasks[0]'\n"
         "\n"
         "\n"
         "# as written in the resource file:\n"
         "\n"
         "tasks:\n"
         "  - task_key: notebook_task\n"
         "    notebook_task:\n"
         "      notebook_path:\n"
         "        ../src/sample_notebook.ipynb",
         "{\n"
         "  \"task_key\": \"notebook_task\",\n"
         "  \"notebook_task\": {\n"
         "    \"notebook_path\":\n"
         "      \"/Workspace/Users/santitham.pro@\n"
         "       kmutt.ac.th/.bundle/m3demo/dev/\n"
         "       files/src/sample_notebook\"\n"
         "  }\n"
         "}\n"
         "\n"
         "# the deployed copy, named before it\n"
         "# exists",
         csize=9.5)

s_challenge(nx(), 3, "Author a second job", "10 min", [
    "Write resources/report.job.yml declaring one job. Its name must come from a new "
    "variable you declare, and it must have exactly one notebook task and no trigger.",
    "Create the notebook the task names, under src/. Two lines are enough.",
    "Assign the new variable a value in the dev target only.",
    "Run databricks bundle validate, then read the resolved job name in -o json.",
], "Validation OK, and a resolved job name carrying the development prefix.")

s_bullets(nx(), "bundle plan",
          "What a deployment would do, reported without doing it", [
    (None, "One line per resource: create, update, delete, or no change."),
    (None, "It compares the resolved configuration against the deployment state, so the first "
           "plan reports every resource as a creation."),
    (None, "--select restricts the plan, and the same flag restricts a deployment."),
    (None, "This is the reviewability property from Part A, available on your own machine "
           "before it is available in a pull request."),
], size=13.5, gap=16,
    note="Module 7 runs bundle plan in a pull request, so that a reviewer reads the effect of a "
         "change rather than only its diff.")

s_bullets(nx(), "bundle deploy, in order",
          "Four stages. A failure message names one of them, and the stage names the cause.", [
    ("1 · Build", "Each entry under artifacts: is built. For the generated project this runs\n"
                  "uv build --wheel, which fails immediately if uv is absent."),
    ("2 · Upload", "The source tree is synchronised to file_path, and build products to\n"
                   "artifact_path. Files excluded by .gitignore are not uploaded."),
    ("3 · Apply", "Resources are created or updated through the API, in dependency order."),
    ("4 · Record", "The deployment state is written, locally and in the workspace."),
], size=12.5, gap=6)

s_shot(nx(), "Where the deployment lives",
       "The workspace after a first deployment to target dev",
       "S46",
       "$ databricks bundle deploy -t dev\n"
       "$ databricks workspace list /Workspace/Users/$USER/.bundle/m3demo/dev",
       "The four directories from Part A, and the job under Workflows carrying the development "
       "prefix. The notebook under files/src is the copy the job executes.")

s_bullets(nx(), "Deployment state",
          "How the CLI knows what it already created", [
    (None, "State records the resources this bundle created, so the next deployment\n"
           "computes a difference rather than recreating everything."),
    (None, "It is held locally under .databricks/bundle/<target>/ and in the workspace\n"
           "under state/. The local directory is git-ignored by the template."),
    (None, "Bundles created with CLI v1.3.0 or later use the direct engine and write\n"
           "resources.json. Earlier bundles use Terraform and write terraform.tfstate."),
    (None, "A bundle carrying both files is refused, reporting the same serial number in\n"
           "terraform and direct states. The migration command is bundle deployment migrate."),
], size=12.5, gap=12,
    note="Deleting the state does not delete the resources. It makes the CLI believe it created "
         "nothing, and the next deployment attempts to create everything again.")

s_challenge(nx(), 4, "Deploy and inspect", "13 min", [
    "Run databricks bundle plan -t dev and read what it intends to do. Then deploy to dev.",
    "In the workspace, find the job. Confirm the name prefix and the paused trigger against "
    "your Challenge 1 prediction.",
    "Navigate to the bundle root path and list its four directories.",
    "State the full workspace path of the notebook this job will execute, and where the same "
    "notebook exists in your project directory.",
], "a deployed job whose name matches your prediction, and both paths of the notebook written down.")

s_shot(nx(), "bundle run",
       "One command replaces the trigger, the polling loop and the output retrieval",
       "S49",
       "$ databricks bundle run sample_job -t dev",
       "The run URL is printed first, then the state transitions, then the terminal result. "
       "Arguments after -- are passed to the job.")

s_split(nx(), "bundle summary and bundle open",
        "Two commands that read deployment state rather than configuration",
        "SUMMARY",
        "$ databricks bundle summary -t dev\n"
        "\n"
        "Lists every resource this bundle has\n"
        "deployed to the selected target, with\n"
        "the workspace id of each.\n"
        "\n"
        "Empty before the first deployment,\n"
        "because state is what it reads.",
        "OPEN",
        "$ databricks bundle open sample_job\n"
        "\n"
        "Opens the resource in a browser,\n"
        "resolving the id from state.\n"
        "\n"
        "Useful when a run has failed and the\n"
        "next step is the run page rather than\n"
        "the terminal.",
        mono_left=True, mono_right=True)

s_challenge(nx(), 5, "Run the job and read the result", "10 min", [
    "Run the job from the command line and record the run URL it prints.",
    "Report the resources this bundle has deployed, and the workspace id of the job.",
    "Open the run and confirm which notebook path the task executed.",
    "Compare that path with the one you wrote down in Challenge 4.",
], "a terminal state obtained by command, and the executed notebook path matching the deployed copy.")

s_shot(nx(), "bundle destroy",
       "Deletes the resources this bundle created, then the files it uploaded",
       "S52",
       "$ databricks bundle destroy -t dev",
       "Two confirmations, resources first and files second. --auto-approve suppresses both, "
       "which is what an automated pipeline passes.")

s_bullets(nx(), "What destroy does not remove",
          "The boundary is the deployment state", [
    (None, "Tables and files written by a job that ran. Destroying the job does not touch what "
           "it produced."),
    (None, "Catalogs and schemas that the bundle referenced rather than declared."),
    (None, "Resources created by hand in the workspace, even where they carry the same name."),
    (None, "Anything deployed under a different target. Destroy acts on one target at a time."),
], size=13.5, gap=16,
    note="Redeploying restores everything the bundle declares, because the declaration and not "
         "the workspace is the source of truth.")

s_challenge(nx(), 6, "Destroy, redeploy, override", "12 min", [
    "Destroy the dev deployment and confirm from the workspace that the job is gone.",
    "Redeploy, passing --var to override the variable you declared in Challenge 3.",
    "Confirm the override in the resolved configuration, and then in the workspace.",
    "Name which of the four sources of a variable value won, and why.",
], "the job restored, the override visible in the workspace, and the precedence rule stated.")

# --- Part D ------------------------------------------------------------------
s_split(nx(), "Your code exists twice",
        "After a deployment, the same notebook exists in two places with different roles",
        "THE SOURCE",
        "my-project/src/sample_notebook.ipynb\n"
        "\n"
        "What you edit. What the repository\n"
        "holds. What a reviewer reads.\n"
        "\n"
        "Editing it changes nothing in the\n"
        "workspace until you deploy.",
        "THE DEPLOYED COPY",
        ".bundle/m3demo/dev/files/src/\n"
        "    sample_notebook\n"
        "\n"
        "What the job executes.\n"
        "\n"
        "Editing it in the workspace changes\n"
        "what runs, and the next deployment\n"
        "overwrites the edit without warning.",
        mono_left=True, mono_right=True,
        note="The task path, read on the validate slide or from the run page, states which copy "
             "a job is running.")

s_table(nx(), "Six failure modes and their signals",
        "In the order you are likely to meet them",
        ["Symptom", "Cause", "Detected"],
        [["cannot configure default credentials", "No profile, and no environment variables",
          "before init writes a file"],
         ["Forbidden on the metastore call", "Credential without Unity Catalog access",
          "before init writes a file"],
         ["uv: command not found", "uv absent; the template builds a wheel with it",
          "build stage of deploy"],
         ["Validation OK, then deploy fails", "The schema is satisfied; the workspace refuses "
          "the operation", "apply stage of deploy"],
         ["The job runs an old version", "The deployed copy was edited, or the source was not "
          "deployed", "not detected; read the task path"],
         ["Deploy to prod refused", "Branch check, or user-specific paths under production mode",
          "before any change is made"]],
        [3.60, 4.20, 2.60], rowh=0.40, size=10, mono_cols=(0,))

s_bullets(nx(), "What this module did not cover",
          "Each of these has a module of its own", [
    ("Authoring jobs, tasks and cluster specifications",
     "Module 4, this afternoon. You read the generated resource file today; tomorrow\n"
     "you write one, with several tasks and a declared cluster."),
    ("Promotion across more than two targets, and variable overrides in earnest",
     "Module 4."),
    ("Diagnosing a deployment that fails in the apply stage",
     "Module 4, with a deliberately sabotaged bundle."),
    ("Tests over the code a bundle deploys, and running these commands on a server",
     "Modules 5 and 7."),
], size=12.5, gap=8)

s_bullets(nx(), "Exit ticket", "Two minutes, on paper, before the break", [
    (None, "Name the four things a deployment writes under the bundle root path."),
    (None, "State what mode: development changes about a job's name and about its trigger."),
    (None, "Name the two sources of a variable value that outrank the declaration's default."),
    (None, "State which copy of a notebook a deployed job executes, and how you would confirm it."),
], size=14, gap=22,
    note="If any of the four is unclear, the answer is in the slides listed on the handout, and "
         "worth settling before Module 4.")

s_end(nx())

TOTAL = i

# stamp the real total into every page number
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and "%%TOTAL%%" in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = r.text.replace("%%TOTAL%%", str(TOTAL))

prs.save("slides.pptx")
print(f"wrote slides.pptx — {TOTAL} slides")
