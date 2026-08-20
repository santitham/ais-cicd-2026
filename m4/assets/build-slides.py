#!/usr/bin/env python3
"""Build Module 4 slides.pptx, matching the Module 1, 2 and 3 decks."""

import copy
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette ---
ORANGE   = C(0xFC, 0x78, 0x00)
INK      = C(0x22, 0x22, 0x24)
BODY     = C(0x3F, 0x3F, 0x46)
MUTED    = C(0x77, 0x77, 0x80)
WHITE    = C(0xFF, 0xFF, 0xFF)
HDR_DARK = C(0x40, 0x40, 0x41)
TERM_BG  = C(0x2E, 0x2E, 0x30)
TERM_TX  = C(0xEC, 0xEC, 0xEE)
GRN_HDR  = C(0x2E, 0x9E, 0x5B)
OUT_BG   = C(0x1C, 0x2C, 0x22)
OUT_TX   = C(0xCC, 0xEE, 0xD8)
OK_BG    = C(0xEC, 0xF7, 0xF0)
NOTE_BG  = C(0xFF, 0xFB, 0xAA)
NOTE_LN  = C(0xF5, 0xD3, 0x28)
RED      = C(0xC0, 0x39, 0x2B)
RED_BG   = C(0xFC, 0xEF, 0xEE)
GREY_BG  = C(0xF4, 0xF4, 0xF6)
GREY_LN  = C(0xD8, 0xD8, 0xDE)

L, W = 2.35, 10.40
FOOT = "Day 2 · Declarative Automation Bundles  ·  Module 4"

prs = Presentation("bx-slide-template.pptx")

# remove the template's example slides
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)

BLANK = prs.slide_masters[0].slide_layouts[6]
TOTAL = 0  # filled in after the build


# ---------------------------------------------------------------- helpers ---
def new():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h, text, size=12, bold=False, color=BODY, font="Arial",
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, line=None):
    tb = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(s, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, In(l), In(t), In(w), In(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def chrome(s, n):
    box(s, L, 6.99, 8.20, 0.30, FOOT, 7.5, False, MUTED)
    box(s, 11.05, 6.99, 1.70, 0.30, f"{n} / %%TOTAL%%", 7.5, False, MUTED,
        align=PP_ALIGN.RIGHT)


def head(s, title, sub=None, size=25):
    box(s, L, 0.36, W, 0.62, title, size, True, INK)
    if sub:
        box(s, L, 1.02, W, 0.40, sub, 11.5, False, MUTED)


# --------------------------------------------------------------- archetypes -
def s_title(n, kicker, title, sub, meta):
    s = new()
    box(s, L, 2.32, W, 0.30, kicker, 11, True, ORANGE)
    box(s, L, 2.72, W, 1.10, title, 33, True, INK)
    box(s, L, 3.98, W, 0.80, sub, 14, False, BODY, line=1.25)
    rect(s, L, 5.10, 1.30, 0.045, ORANGE)
    box(s, L, 5.44, W, 0.30, meta, 11, False, MUTED)
    box(s, L, 6.62, W, 0.28, "DataOps Mastery — Advanced CI/CD and Pipeline Automation",
        9.5, False, MUTED)
    chrome(s, n)


def s_part(n, kicker, title, sub):
    s = new()
    rect(s, L, 3.02, W, 0.04, ORANGE)
    box(s, L, 2.56, W, 0.34, kicker, 11, True, ORANGE)
    box(s, L, 3.24, W, 1.30, title, 33, True, INK)
    box(s, L, 4.58, W, 0.80, sub, 13, False, MUTED, line=1.25)
    chrome(s, n)


def s_bullets(n, title, sub, items, note=None, size=13.5, gap=11):
    """items: list of (bold_lead_or_None, text)"""
    s = new()
    head(s, title, sub)
    t = 1.68
    for lead, text in items:
        rect(s, L, t + 0.09, 0.075, 0.075, ORANGE, shape=MSO_SHAPE.OVAL)
        if lead:
            box(s, L + 0.28, t - 0.03, W - 0.28, 0.28, lead, size, True, INK)
            tb = box(s, L + 0.28, t + 0.25, W - 0.28, 0.60, text, size - 1.5,
                     False, BODY, line=1.22)
            t += 0.30 + 0.24 * (1 + text.count("\n")) + gap / 100.0 + 0.16
        else:
            box(s, L + 0.28, t - 0.03, W - 0.28, 0.60, text, size, False, BODY,
                line=1.22)
            t += 0.24 * (1 + text.count("\n")) + gap / 100.0 + 0.10
    if note:
        y = max(t + 0.14, 5.95)
        rect(s, L, y, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_runsee(n, title, sub, run, see, run_label="YOU RUN", see_label="YOU SEE",
             see_ok=True, csize=9.5):
    s = new()
    head(s, title, sub)
    hdr2 = GRN_HDR if see_ok else C(0x9E, 0x3A, 0x2E)
    bg2 = OUT_BG if see_ok else C(0x2E, 0x1C, 0x1C)
    tx2 = OUT_TX if see_ok else C(0xF2, 0xD2, 0xCE)
    rect(s, L, 1.62, 5.03, 0.34, HDR_DARK)
    box(s, L + 0.14, 1.62, 4.75, 0.34, run_label, 9, True, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, L, 1.96, 5.03, 4.92, TERM_BG)
    box(s, L + 0.16, 2.08, 4.71, 4.68, run, csize, False, TERM_TX, "Consolas",
        line=1.24)
    rect(s, 7.72, 1.62, 5.03, 0.34, hdr2)
    box(s, 7.86, 1.62, 4.75, 0.34, see_label, 9, True, WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 7.72, 1.96, 5.03, 4.92, bg2)
    box(s, 7.88, 2.08, 4.71, 4.68, see, csize, False, tx2, "Consolas", line=1.24)
    chrome(s, n)


def s_code(n, title, sub, code, note=None, csize=11, ch=4.60):
    s = new()
    head(s, title, sub)
    rect(s, L, 1.62, W, ch, TERM_BG)
    box(s, L + 0.16, 1.74, W - 0.32, ch - 0.24, code, csize, False, TERM_TX,
        "Consolas", line=1.26)
    if note:
        y = 1.62 + ch + 0.16
        rect(s, L, y, W, 0.56, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.56, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    chrome(s, n)


NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # "No Style, No Grid"


def _cell_text(cell, text, size, bold, color, font, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.16
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def s_table(n, title, sub, headers, rows, widths, note=None, size=11.5,
            rowh=0.42, mono_cols=()):
    """A real PowerPoint table: editable cells, resizable columns."""
    s = new()
    head(s, title, sub)

    hdr_h = 0.40
    row_hs = [rowh * (1 + max(c.count("\n") for c in row)) for row in rows]
    top = 1.66

    gf = s.shapes.add_table(len(rows) + 1, len(headers),
                            In(L), In(top), In(W), In(hdr_h + sum(row_hs)))
    tbl = gf.table

    # Suppress the built-in blue table style so our own fills show through.
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for sid in tblPr.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"):
        sid.text = NO_STYLE

    for ci, w in enumerate(widths):
        tbl.columns[ci].width = In(w)
    tbl.rows[0].height = In(hdr_h)
    for ri, h in enumerate(row_hs, start=1):
        tbl.rows[ri].height = In(h)

    def style_cell(cell, fill):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = In(0.14)
        cell.margin_right = In(0.10)
        cell.margin_top = In(0.04)
        cell.margin_bottom = In(0.04)

    for ci, htxt in enumerate(headers):
        c = tbl.cell(0, ci)
        style_cell(c, HDR_DARK)
        _cell_text(c, htxt, 10.5, True, WHITE, "Arial")

    for ri, row in enumerate(rows):
        for ci, cell_txt in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            style_cell(c, GREY_BG if ri % 2 == 0 else WHITE)
            col = INK if ci == 0 else BODY
            bold = ci == 0
            if cell_txt.startswith("*"):
                cell_txt, col, bold = cell_txt[1:], ORANGE, True
            _cell_text(c, cell_txt, size, bold, col,
                       "Consolas" if ci in mono_cols else "Arial")

    if note:
        y = min(top + hdr_h + sum(row_hs) + 0.20, 6.14)
        rect(s, L, y, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, y, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


def s_challenge(n, num, title, minutes, steps, success):
    s = new()
    rect(s, L, 0.30, W, 1.24, HDR_DARK)
    rect(s, L, 1.54, W, 0.05, ORANGE)
    box(s, L + 0.24, 0.44, 4.00, 0.30, "LAB", 11, True, ORANGE)
    box(s, L + 0.24, 0.76, 7.80, 0.66, f"Challenge {num} — {title}", 21, True, WHITE)
    rect(s, 10.85, 0.46, 1.60, 0.38, ORANGE)
    box(s, 10.85, 0.46, 1.60, 0.38, minutes, 9.5, True, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    t = 1.96
    step = min(0.62, (4.90 - 0.10) / max(len(steps), 1))
    for i, txt in enumerate(steps, 1):
        rect(s, L + 0.10, t, 0.34, 0.34, ORANGE, shape=MSO_SHAPE.OVAL)
        box(s, L + 0.10, t, 0.34, 0.34, str(i), 9, True, WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        box(s, L + 0.58, t - 0.04, 9.72, step - 0.06, txt, 12, False, BODY,
            line=1.20)
        t += step
    y = max(t + 0.12, 6.10)
    rect(s, L, y, W, 0.50, OK_BG, GRN_HDR)
    box(s, L + 0.18, y, W - 0.36, 0.50, "Success: " + success, 11.5, False, INK,
        anchor=MSO_ANCHOR.MIDDLE)
    chrome(s, n)


def s_shot(n, title, sub, ref, commands, what):
    """Screenshot placeholder. `ref` matches assets/capture-screenshots.md."""
    s = new()
    head(s, title, sub)
    rect(s, L, 1.62, W, 0.86, TERM_BG)
    box(s, L + 0.16, 1.72, W - 0.32, 0.68, commands, 10.5, False, TERM_TX,
        "Consolas", line=1.24)
    ph = rect(s, L, 2.64, W, 3.42, C(0xFA, 0xFA, 0xFB), GREY_LN)
    box(s, L, 3.66, W, 0.34, f"SCREENSHOT {ref}", 15, True, MUTED,
        align=PP_ALIGN.CENTER)
    box(s, L, 4.06, W, 0.34, "capture on the Ubuntu server — assets/capture-screenshots.md",
        10.5, False, MUTED, align=PP_ALIGN.CENTER)
    rect(s, L, 6.20, W, 0.56, NOTE_BG, NOTE_LN)
    box(s, L + 0.18, 6.20, W - 0.36, 0.56, what, 11.5, False, INK,
        anchor=MSO_ANCHOR.MIDDLE, line=1.16)
    chrome(s, n)


def s_split(n, title, sub, left_title, left_body, right_title, right_body,
            note=None, mono_left=False, mono_right=False):
    s = new()
    head(s, title, sub)
    for x, ttl, bdy, mono in ((L, left_title, left_body, mono_left),
                              (7.72, right_title, right_body, mono_right)):
        rect(s, x, 1.62, 5.03, 0.40, HDR_DARK)
        box(s, x + 0.16, 1.62, 4.71, 0.40, ttl, 10.5, True, WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x, 2.02, 5.03, 3.90, GREY_BG, GREY_LN)
        box(s, x + 0.18, 2.18, 4.67, 3.58, bdy,
            10 if mono else 12.5, False, BODY,
            "Consolas" if mono else "Arial", line=1.28)
    if note:
        rect(s, L, 6.14, W, 0.62, NOTE_BG, NOTE_LN)
        box(s, L + 0.18, 6.14, W - 0.36, 0.62, note, 11.5, False, INK,
            anchor=MSO_ANCHOR.MIDDLE, line=1.18)
    chrome(s, n)


# ================================================================== CONTENT ==
i = 0
def nx():
    global i
    i += 1
    return i


def s_end(n):
    s = new()
    box(s, L, 2.70, W, 0.30, "END OF MODULE 4", 11, True, ORANGE)
    box(s, L, 3.10, W, 0.90, "Next: Testing pipelines", 30, True, INK)
    box(s, L, 4.14, 8.60, 0.90,
        "Day 2 verified declarations. Module 3 confirmed the workspace received what the "
        "file said, and Module 4 confirmed a task failed when its code was broken. "
        "Neither confirmed that a task which succeeded computed the right answer. "
        "Module 5 does.",
        13.5, False, BODY, line=1.28)
    rect(s, L, 5.42, 1.30, 0.045, ORANGE)
    box(s, L, 5.74, W, 0.30, "santitham.pro@kmutt.ac.th", 12, False, MUTED)
    chrome(s, n)


# --- front matter ------------------------------------------------------------
s_title(nx(), "MODULE 4", "Build, deploy, run, iterate",
        "Writing the resource files you read this morning: jobs with several tasks, "
        "cluster specifications, and one bundle promoted across targets.",
        "Day 2 · afternoon · 3 hours")

s_bullets(nx(), "Module overview", "By the end of this module, you can", [
    (None, "Author a job resource from written requirements: tasks, payloads and compute."),
    (None, "Declare a job cluster, bind tasks to it, and say what each field controls."),
    (None, "Express a dependency graph, and predict each task's state when one fails."),
    (None, "Import an existing workspace job and correct the defects it arrives with."),
    (None, "Deploy one bundle to two targets, overriding only what differs."),
    (None, "Diagnose a failed deployment by identifying the stage that reported it."),
], size=14, gap=20)

s_table(nx(), "How this module runs",
        "Three hours: 105 minutes of instruction, 60 at your own keyboard",
        ["", "Topic", "Teach", "You work"],
        [["Part A", "Declaring a job: tasks, payloads and paths", "20 min", "—"],
         ["Part B", "Compute: job clusters and how a task claims one", "16 min", "8 min"],
         ["Part C", "Importing a job that already exists", "10 min", "10 min"],
         ["Part D", "Multi-task pipelines and dependencies", "16 min", "12 min"],
         ["Break", "", "15 min", ""],
         ["Part E", "Targets, overrides and variable precedence", "17 min", "10 min"],
         ["Part F", "The loop, and the diagnosis recipe", "15 min", "20 min"],
         ["Part G", "What comes next, exit ticket", "5 min", "—"],
         ["*Total", "*3 hours", "*105 min", "*60 min"]],
        [1.90, 4.70, 1.90, 1.90], rowh=0.32, size=11,
        note="Challenges 3 to 6 each build on the bundle the previous one left behind. If a "
             "challenge does not complete, take the solution directory and move forward.")

s_table(nx(), "Where this sits",
        "This morning you read a file a template wrote. This afternoon you write it.",
        ["", "Module 3", "Module 4"],
        [["The file", "generated by bundle init", "*written by you"],
         ["Resources", "one job and one pipeline, given", "*authored from requirements"],
         ["Tasks", "three, read", "*three, with declared dependencies"],
         ["Targets", "two, as generated", "*three, with deliberate overrides"],
         ["Failure", "the lifecycle succeeded", "*a deployment is diagnosed and repaired"]],
        [2.10, 4.05, 4.25], rowh=0.44, size=11.5)

# --- Part A ------------------------------------------------------------------
s_part(nx(), "PART A", "Declaring a job",
       "The keys that identify a job and name the work it performs. One key per slide, "
       "starting from a bundle that has no resources at all.")

s_code(nx(), "The starting point", "bundle init default-minimal — a bundle with nowhere for a job to go yet",
       "$ databricks bundle init default-minimal\n"
       "$ tree -a -I '.databricks|.git'\n"
       ".\n"
       "|-- databricks.yml        bundle, include, variables, dev and prod targets\n"
       "|-- pyproject.toml        package metadata\n"
       "|-- tests/conftest.py\n"
       "|-- fixtures/.gitkeep\n"
       "|-- AGENTS.md  CLAUDE.md  README.md\n"
       "`-- .vscode/\n"
       "\n"
       "$ databricks bundle validate\n"
       "Validation OK!",
       note="databricks.yml declares include: [resources/*.yml], and the directory does not "
            "exist. A glob that matches nothing is not an error. Everything you write today "
            "goes into that directory.", csize=11, ch=3.10)

s_bullets(nx(), "Three names above every job",
          "resources -> the resource type -> a key you choose", [
    ("resources",
     "The top-level key you met this morning. Everything the workspace should contain\n"
     "is declared beneath it."),
    ("jobs",
     "The resource type. Pluralised, and fixed by the schema: jobs, pipelines,\n"
     "clusters, schemas, volumes, model_serving_endpoints, dashboards."),
    ("a key you choose",
     "The identifier you will pass to bundle run and reference from other resources.\n"
     "It is not the job's name, and it never appears in the workspace."),
], size=13.5, gap=16,
   note="This morning the key sample_job was given to you. From here on you are choosing it, "
        "so choose it to be typed: bundle run is how you will start this job all afternoon.")

s_code(nx(), "The smallest job", "Five lines. Every slide in Parts A and B adds or replaces exactly one of them.",
       "# resources/ingest.job.yml\n"
       "resources:\n"
       "  jobs:\n"
       "    ingest:                       <-- the key you chose\n"
       "      name: sales_ingest       <-- what the workspace will display\n"
       "      tasks:\n"
       "        - task_key: land\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/land.py",
       note="This validates and deploys. It creates a job with one task, and the task has no "
            "compute declared, which Part B fixes.", csize=12.5, ch=2.72)

s_bullets(nx(), "task_key is an identity, not a label",
          "What refers to it, and what does not check it", [
    ("Three things refer to a task by its key",
     "depends_on, in Part D. bundle run --only, when you rerun one task of a deployed\n"
     "job. And the run history, which reports state per key."),
    ("Two tasks may carry the same key without complaint",
     "bundle validate reports Validation OK! for a job with two tasks keyed clean.\n"
     "The Jobs API rejects it at deploy. Uniqueness is enforced there, not locally."),
], size=13.5, gap=20,
   note="This is the first instance of a pattern that runs through the whole module: the CLI "
        "checks the schema and the references it owns, and the workspace checks the rest.")

s_table(nx(), "One payload key per task",
        "The key that names the work. A task declares exactly one of these.",
        ["Payload key", "Runs", "Names"],
        [["notebook_task", "a notebook in the workspace", "notebook_path, base_parameters"],
         ["spark_python_task", "a Python file, with Spark available", "python_file, parameters"],
         ["python_wheel_task", "an entry point in a built wheel", "package_name, entry_point"],
         ["sql_task", "a query or file on a SQL warehouse", "warehouse_id, then query or file"],
         ["run_job_task", "another job", "job_id"],
         ["condition_task", "a comparison that gates what follows", "op, left, right"]],
        [2.75, 4.15, 3.50], rowh=0.42, size=11, mono_cols=(0, 2),
        note="This course uses the first two. python_wheel_task is what the default-python "
             "template you saw this morning declares, which is why that template needs uv.")

s_code(nx(), "notebook_task", "The path is relative to the file that declares it, not to the bundle root",
       "# resources/ingest.job.yml            the file is in resources/\n"
       "      tasks:\n"
       "        - task_key: land\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/land.py\n"
       "            base_parameters:\n"
       "              catalog: ${var.catalog}\n"
       "              schema: ${var.schema}",
       note="base_parameters are read inside the notebook with dbutils.widgets.get. This is how "
            "a notebook receives the catalog for the target it is running in, rather than "
            "hard-coding one.", csize=12.5, ch=2.48)

s_code(nx(), "spark_python_task", "The same task, with the payload replaced",
       "      tasks:\n"
       "        - task_key: land\n"
       "          spark_python_task:\n"
       "            python_file: ../src/land.py\n"
       "            parameters:\n"
       "              - --catalog\n"
       "              - ${var.catalog}\n"
       "              - --lookback-days\n"
       "              - \"7\"",
       note="parameters is a sequence of strings, passed to the file as command-line arguments "
            "and read with argparse. Note the quotes on \"7\": Module 2's YAML rules still apply, "
            "and an unquoted 7 is an integer where a string is required.", csize=12.5, ch=2.72)

s_table(nx(), "A local path becomes a workspace path",
        "Resolved before any deployment, and the rewriting differs by payload key",
        ["Declared in resources/ingest.job.yml", "Resolved for target dev"],
        [["notebook_task:\n  notebook_path: ../src/land.py",
          ".../.bundle/m4demo/dev/files/src/land"],
         ["spark_python_task:\n  python_file: ../src/land.py",
          ".../.bundle/m4demo/dev/files/src/land.py"]],
        [4.90, 5.50], rowh=0.34, size=10.5, mono_cols=(0, 1),
        note="The notebook loses its extension because a notebook is a workspace object, not a "
             "file. The Python file keeps it because it stays a file. Read either value with "
             "bundle validate -o json before you deploy.")

s_bullets(nx(), "What makes a .py file a notebook",
          "The one rule in Part A enforced by reading the file rather than its path", [
    ("The first line must be the notebook marker",
     "# Databricks notebook source\n"
     "Nothing else distinguishes a notebook from a Python file on disk. The marker is\n"
     "what the workspace import path looks for."),
    ("Without it, the same declaration is rejected",
     "Error: expected a notebook for \"resources.jobs.ingest.tasks[1]\n"
     ".notebook_task.notebook_path\" but got a file: file at .../src/land.py\n"
     "is not a notebook"),
], size=13.5, gap=18,
   note="A file you wrote by hand will not have the marker unless you add it. A file exported "
        "from the workspace will. This is the most common failure in the next hour.")

s_table(nx(), "Keys that belong on the job, not on a task",
        "Four you will use this afternoon",
        ["Key", "Controls", "Note"],
        [["max_concurrent_runs", "how many runs may overlap", "development mode sets 4"],
         ["timeout_seconds", "the whole job's limit, 0 for none", "tasks have their own"],
         ["tags", "key/value labels in the workspace", "development mode adds one"],
         ["email_notifications", "on_start, on_success, on_failure", "each is a list of addresses"]],
        [3.20, 4.40, 2.80], rowh=0.44, size=11, mono_cols=(0,),
        note="Anything development mode supplies, you saw this morning as a preset. Anything "
             "else, including who gets the failure mail, has to be declared here.")

# --- Part B ------------------------------------------------------------------
s_part(nx(), "PART B", "Compute for a task",
       "Three ways a task obtains a cluster, why this course declares one of them, "
       "and what each field of a cluster specification fixes.")

s_table(nx(), "Three ways a task obtains compute",
        "The declaration, and the consequence of each",
        ["Declared on the task", "Compute", "Started"],
        [["job_cluster_key: main", "a cluster this job creates", "at run time, terminated after"],
         ["existing_cluster_id: 0815-...", "an interactive cluster already there", "it is already running"],
         ["environment_key: default", "serverless", "managed by the platform"]],
        [3.60, 3.70, 3.10], rowh=0.46, size=11, mono_cols=(0,),
        note="existing_cluster_id names one cluster in one workspace, so a bundle carrying it "
             "cannot be promoted. That is the reason this course declares job clusters.")

s_code(nx(), "job_clusters is declared on the job", "A sequence. Each entry requires exactly two keys.",
       "resources:\n"
       "  jobs:\n"
       "    ingest:\n"
       "      name: sales_ingest\n"
       "      job_clusters:\n"
       "        - job_cluster_key: main       <-- the identity of this entry\n"
       "          new_cluster:                <-- the specification\n"
       "            spark_version: 16.4.x-scala2.12\n"
       "            node_type_id: Standard_D3_v2\n"
       "            data_security_mode: DATA_SECURITY_MODE_AUTO\n"
       "            autoscale:\n"
       "              min_workers: 1\n"
       "              max_workers: 4\n"
       "      tasks:\n"
       "        - task_key: land\n"
       "          job_cluster_key: main       <-- the reference to it\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/land.py",
       csize=11, ch=4.52)

s_bullets(nx(), "job_cluster_key appears twice",
          "Once as an identity, once for each task that claims it", [
    ("On the job_clusters entry, it is the name of the specification",
     "Nothing else in the file gives that specification a handle."),
    ("On a task, it is a claim on that specification",
     "Two tasks naming the same key run on one cluster, one after the other. Two tasks\n"
     "naming different keys run on two clusters, and may run at the same time."),
    ("A cluster no task claims is still created",
     "It appears in the deployed job definition and starts nothing. Part E shows the\n"
     "override mistake that produces one."),
], size=13.5, gap=13)

s_table(nx(), "The fields of new_cluster",
        "Four fields fix everything this course depends on",
        ["Field", "Value used here", "Fixes"],
        [["spark_version", "16.4.x-scala2.12", "the runtime, and so Spark and Python"],
         ["node_type_id", "Standard_D3_v2", "the instance type of driver and workers"],
         ["data_security_mode", "DATA_SECURITY_MODE_AUTO", "the access mode, and what UC permits"],
         ["autoscale / num_workers", "see the next slide", "the worker count"]],
        [3.10, 3.60, 3.70], rowh=0.44, size=10.5, mono_cols=(0, 1),
        note="16.4 is the long-term support runtime the Databricks templates currently pin, "
             "supported to May 2028. Documentation examples still showing 13.3.x-scala2.12 name "
             "a runtime that goes out of support in August 2026.")

s_split(nx(), "num_workers or autoscale, not both",
        "Two ways to state the worker count",
        "A FIXED COUNT",
        "num_workers: 2\n"
        "\n"
        "Two workers for the whole run.\n"
        "Predictable cost. Predictable\n"
        "duration. Nothing adapts if the\n"
        "input is larger than usual.",
        "A RANGE",
        "autoscale:\n"
        "  min_workers: 1\n"
        "  max_workers: 4\n"
        "\n"
        "Starts at one and adds workers\n"
        "while there is queued work.\n"
        "max_workers must exceed\n"
        "min_workers.",
        mono_left=True, mono_right=True,
        note="Declaring both passes validation and resolves with both present. The conflict is "
             "not reported until the Jobs API receives it at deploy.")

s_split(nx(), "One cluster for the job, or one for each task",
        "The same three tasks, declared two ways",
        "ONE ENTRY, THREE TASKS CLAIM IT",
        "One cluster start, once, for the\n"
        "whole run.\n\n"
        "The three tasks run in sequence on\n"
        "shared compute even where the graph\n"
        "would allow two of them in parallel.\n\n"
        "One size has to suit all three.\n\n"
        "Cheapest, and slowest.",
        "THREE ENTRIES, ONE EACH",
        "Three cluster starts. Each is a few\n"
        "minutes of the run.\n\n"
        "Tasks the graph allows in parallel\n"
        "actually run in parallel.\n\n"
        "Each task is sized for its own work:\n"
        "the features task gets the workers,\n"
        "the notify task does not.\n\n"
        "Faster, and dearer.",
        note="The choice is between start-up cost and per-task sizing, and it is made in the "
             "file. Nothing outside the declaration adjusts it.")

s_runsee(nx(), "What validate reports when the binding is wrong",
         "A task claims mian; the entry is keyed main",
         "$ databricks bundle validate\n"
         "\n"
         "\n"
         "# resources/ingest.job.yml, extract\n"
         "      job_clusters:\n"
         "        - job_cluster_key: main\n"
         "          new_cluster:\n"
         "            ...\n"
         "      tasks:\n"
         "        - task_key: land\n"
         "          job_cluster_key: mian\n"
         "            ^ one transposition",
         "Warning: job_cluster_key mian is\n"
         "not defined\n"
         "  at resources.jobs.ingest\n"
         "     .tasks[1].job_cluster_key\n"
         "  in resources/ingest.job.yml:42:28\n"
         "\n"
         "Name: m4demo\n"
         "Target: dev\n"
         "...\n"
         "\n"
         "Found 1 warning",
         see_ok=False, csize=10)

s_bullets(nx(), "Read the whole message, not the last line",
          "The warning above carries three separate pieces of information", [
    ("It is a warning, not an error",
     "The command reports Found 1 warning and exits successfully. A participant who\n"
     "reads only the last line will deploy a job whose task has no compute."),
    ("at resources.jobs.ingest.tasks[1].job_cluster_key",
     "The configuration path: the address in the resolved configuration, which is what\n"
     "you query with bundle validate -o json."),
    ("in resources/ingest.job.yml:42:28",
     "The file, the line and the column. This is where you edit."),
], size=13.5, gap=11,
   note="land is the FIRST task in the file, and the path says tasks[1], because the resolved "
        "configuration sorts a job's tasks by task_key: features, land, notify, publish, "
        "quality. Counting down the file opens the wrong task. Edit by position, not by index.")

s_challenge(nx(), 1, "Author a job on a job cluster", "8 MIN", [
    "Read assets/challenge1-requirements.md. Do not open any existing resource file.",
    "Write resources/ingest.job.yml from the requirements: one job, one notebook task, "
    "one job cluster that you declare.",
    "Create the notebook the task names, with the marker line as its first line.",
    "Before you validate, write down the workspace path you expect notebook_path to "
    "resolve to.",
    "Run databricks bundle validate, then read the resolved path out of the JSON output "
    "and compare it with your prediction.",
], "Validation OK!, and a resolved notebook_path matching what you wrote down.")

# --- Part C ------------------------------------------------------------------
s_part(nx(), "PART C", "Importing a job that already exists",
       "Most jobs are not authored from requirements. They exist in the workspace, "
       "built through the interface, and their definition has to be recovered.")

s_bullets(nx(), "bundle generate job",
          "One command reads a job out of the workspace and writes it as bundle source", [
    ("databricks bundle generate job --existing-job-id <id> --key <key>",
     "Reads the job through the Jobs API, writes a resource file, and downloads every\n"
     "notebook the job names."),
    ("Two destinations, both with defaults",
     "--config-dir, default resources/, receives the YAML. --source-dir, default src/,\n"
     "receives the notebooks. --key sets the resource key; without it the CLI picks one."),
    ("Only notebook tasks are supported",
     "A job whose tasks are wheel or SQL tasks is not generated. Add\n"
     "--download-spark-python-files to bring down the files a spark_python_task names."),
], size=13, gap=13,
   note="This is the command that ends the argument about whether to rebuild a job by hand. "
        "You do not have to: you have to correct what it gives you.")

s_code(nx(), "What it produces", "Generated from a one-task job built in the interface on an interactive cluster",
       "resources:\n"
       "  jobs:\n"
       "    hello_world:\n"
       "      name: hello_world_ui\n"
       "      tasks:\n"
       "        - task_key: hello\n"
       "          existing_cluster_id: 0815-123456-abcdefgh\n"
       "          email_notifications: {}\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/hello_world.py\n"
       "            source: WORKSPACE\n"
       "          run_if: ALL_SUCCESS\n"
       "          timeout_seconds: 0\n"
       "          webhook_notifications: {}\n"
       "      email_notifications: {}\n"
       "      max_concurrent_runs: 1\n"
       "      queue:\n"
       "        enabled: true\n"
       "      timeout_seconds: 0\n"
       "      webhook_notifications: {}",
       csize=10.5, ch=4.74)

s_table(nx(), "Three defects, and what reports them",
        "Every one of these passes bundle validate",
        ["In the generated file", "Why it is wrong", "Reported by"],
        [["existing_cluster_id: 0815-...",
          "names a cluster in one workspace;\nno other target has it",
          "nothing, until you deploy\nto a second target"],
         ["source: WORKSPACE",
          "the file is bundle source now,\nsynchronised with the bundle",
          "nothing"],
         ["nine keys carrying defaults\nor empty mappings",
          "no information, and they hide the\nnine lines that carry some",
          "nothing"]],
        [3.30, 3.70, 3.40], rowh=0.30, size=10.5, mono_cols=(0,),
        note="The generated file is a starting point, not an output. Reading it is the fastest "
             "way to learn the shape of a job, which is why this part follows Parts A and B "
             "rather than replacing them.")

s_bullets(nx(), "Adopting the job instead of creating a second one",
          "What happens on the first deployment, and how to change it", [
    ("Deployed as it stands, the generated resource creates a new job",
     "The bundle has no record of the original, so it creates one and leaves the\n"
     "hand-built job in place. Two jobs with nearly the same name, one of them\n"
     "unmanaged."),
    ("bundle deployment bind <key> <job-id>",
     "Records the existing job in the deployment state under your resource key. The\n"
     "next deploy updates that job rather than creating another."),
    ("bundle deployment unbind <key>",
     "Removes the record without deleting the job, which is what you do before handing\n"
     "a job back to someone who manages it in the interface."),
], size=13, gap=13,
   note="Binding is also the answer to the question the room asks in Module 3: what do I do "
        "about everything already in the workspace.")

s_challenge(nx(), 2, "Port a job and repair it", "10 MIN", [
    "Take the job id from the instructor and run bundle generate job with --key ported.",
    "Read the generated file and mark the three defects from the previous slide before "
    "you change anything.",
    "Replace existing_cluster_id with a job cluster you declare, remove the keys that "
    "carry only defaults, and drop source: WORKSPACE.",
    "Validate, then confirm from the JSON output that the resolved notebook_path is "
    "inside the bundle's files root rather than your home directory.",
], "A resource file of under fifteen lines, all of them carrying information, that validates.")

# --- Part D ------------------------------------------------------------------
s_part(nx(), "PART D", "Multi-task pipelines",
       "How one task waits for another, what graph the declaration produces, and what "
       "happens to everything downstream when a task fails.")

s_bullets(nx(), "depends_on",
          "A sequence of mappings, each naming one task_key", [
    ("A task with depends_on does not start until its dependencies are terminal",
     "Terminal means finished, in any state. Whether the task then runs is decided by\n"
     "run_if, three slides from here."),
    ("A task without depends_on starts when the run starts",
     "Several such tasks start together. This is where parallelism comes from; nothing\n"
     "else in the file requests it."),
    ("The dependency is declared downstream",
     "land does not say what follows it; features says what it waits for. The order of\n"
     "entries in the tasks sequence therefore carries no meaning at all."),
], size=13, gap=13)

s_code(nx(), "A three-task chain", "land, then features, then publish — one cluster, claimed three times",
       "      tasks:\n"
       "        - task_key: land\n"
       "          job_cluster_key: main\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/land.py\n"
       "\n"
       "        - task_key: features\n"
       "          depends_on:\n"
       "            - task_key: land\n"
       "          job_cluster_key: main\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/features.py\n"
       "\n"
       "        - task_key: publish\n"
       "          depends_on:\n"
       "            - task_key: features\n"
       "          job_cluster_key: main\n"
       "          notebook_task:\n"
       "            notebook_path: ../src/publish.py",
       csize=11, ch=4.72)

s_split(nx(), "The graph a declaration produces",
        "Three shapes, and the depends_on that produces each",
        "FAN-OUT: TWO TASKS, ONE PARENT",
        "features:  depends_on: [land]\n"
        "quality:   depends_on: [land]\n"
        "\n"
        "         land\n"
        "        /    \\\n"
        "  features    quality\n"
        "\n"
        "Both start when land succeeds, and\n"
        "run at the same time if each has its\n"
        "own cluster.",
        "FAN-IN: ONE TASK, TWO PARENTS",
        "publish:\n"
        "  depends_on:\n"
        "    - task_key: features\n"
        "    - task_key: quality\n"
        "\n"
        "  features    quality\n"
        "        \\    /\n"
        "        publish\n"
        "\n"
        "publish waits for both.",
        mono_left=True, mono_right=True,
        note="A chain is the degenerate case of both. Nothing in the file draws the graph; the "
             "graph is what the set of depends_on entries implies.")

s_table(nx(), "run_if",
        "Six values, and the CLI prints all six when you give it a seventh",
        ["Value", "The task runs when its dependencies", "Used for"],
        [["*ALL_SUCCESS", "*all succeeded — the default", "*the ordinary case"],
         ["AT_LEAST_ONE_SUCCESS", "at least one succeeded", "partial-input work"],
         ["NONE_FAILED", "none failed, skipped or not", "optional upstream steps"],
         ["ALL_DONE", "all reached a terminal state", "cleanup, notification"],
         ["AT_LEAST_ONE_FAILED", "at least one failed", "compensating work"],
         ["ALL_FAILED", "all failed", "a fallback path"]],
        [3.30, 4.35, 2.75], rowh=0.38, size=10.5, mono_cols=(0,),
        note="ALL_DONE is what a notification task declares, so that it runs whether or not the "
             "pipeline succeeded. ALL_DONE_OR_SKIPPED and NONE_FAILED_OR_SKIPPED do not exist.")

s_table(nx(), "What a failure does to the tasks below it",
        "The states a run reports, and which of them the code produced",
        ["State", "Meaning", "Produced by"],
        [["SUCCESS", "the task finished", "the code"],
         ["FAILED", "the task raised, or exited non-zero", "the code"],
         ["TIMEDOUT", "the task exceeded timeout_seconds", "the code, or the cluster"],
         ["UPSTREAM_FAILED", "a dependency failed, so this never ran", "the platform"],
         ["SKIPPED", "run_if was not satisfied", "the platform"],
         ["CANCELED", "the run was stopped", "a person, or a concurrency limit"]],
        [3.05, 4.55, 2.80], rowh=0.36, size=10.5, mono_cols=(0,),
        note="A downstream task of a failed task is not reported as a failure. The run history "
             "distinguishes the task that broke from the tasks that were prevented, which is "
             "what lets you find the first cause rather than the last symptom.")

s_bullets(nx(), "Retries", "Declared per task, and they re-run only that task", [
    ("max_retries",
     "How many further attempts after the first failure. -1 means retry indefinitely,\n"
     "which is a decision, not a default."),
    ("min_retry_interval_millis",
     "The wait between the end of one attempt and the start of the next. Without it,\n"
     "a task that fails on a rate limit retries straight into the same rate limit."),
    ("retry_on_timeout",
     "Whether a TIMEDOUT attempt counts as retriable. False by default, because a task\n"
     "that ran out of time will usually do so again."),
], size=13.5, gap=15,
   note="A retry re-runs the task on the same cluster and does not re-run anything above it. "
        "A task that is not idempotent must not be given retries.")

s_bullets(nx(), "Three graph faults that validate accepts",
          "All three return Validation OK!, and all three are rejected at deploy", [
    ("depends_on naming a task that does not exist",
     "A transposed key in a dependency reads as a dependency on nothing."),
    ("A cycle",
     "land depends on publish, publish depends on features, features depends on land.\n"
     "No task can start."),
    ("Two tasks sharing a task_key",
     "The second is not a second task; the definition is ambiguous."),
], size=13.5, gap=17,
   note="Validation checks the schema and resolves the references it owns. It does not evaluate "
        "the graph, so the resource stage of deploy is the first place these appear. That is "
        "the distinction the recipe in Part F is built on.")

s_challenge(nx(), 3, "Build the pipeline", "12 MIN", [
    "Read assets/challenge3-requirements.md and draw the graph it describes before "
    "writing anything.",
    "Declare four tasks: one root, two that fan out from it, and one that fans in from "
    "both.",
    "Give the fan-out tasks their own job clusters so they can actually run in "
    "parallel, and the others the shared one.",
    "Add a fifth task that notifies, with run_if: ALL_DONE and no cluster of its own.",
    "Give one retry to the task most likely to fail, with an interval, and say in one "
    "line why that task and not another.",
    "Validate, then check your drawn graph against the depends_on entries you wrote.",
], "Validation OK!, and a drawing that matches the file rather than your intention.")

# --- Part E ------------------------------------------------------------------
s_part(nx(), "PART E", "Targets and overrides",
       "One declaration, several destinations. What a second environment actually "
       "differs by, and the two mechanisms for expressing it.")

s_table(nx(), "What a second environment differs by",
        "Four properties. A target block that contains anything else is doing too much.",
        ["Property", "Expressed by", "dev to staging"],
        [["the workspace", "targets.<t>.workspace.host", "may be the same host"],
         ["where data is written", "a variable, per target", "catalog and schema differ"],
         ["the size of the compute", "a cluster override, or a complex variable", "1-4 workers to 2-8"],
         ["whether schedules fire", "mode, or trigger.pause_status", "paused to unpaused"]],
        [3.05, 4.35, 3.00], rowh=0.46, size=10.5,
        note="Everything else is identical, and identical is the point. A target that overrides "
             "the task graph is two pipelines maintained in one file.")

s_code(nx(), "Adding a staging target", "mode: production, an explicit root path, and the two variables reassigned",
       "targets:\n"
       "  dev:\n"
       "    mode: development\n"
       "    default: true\n"
       "    variables:\n"
       "      catalog: training\n"
       "      schema: ${workspace.current_user.short_name}\n"
       "\n"
       "  staging:\n"
       "    mode: production\n"
       "    workspace:\n"
       "      root_path: /Workspace/Users/${workspace.current_user.userName}/\\\n"
       "                 .bundle/${bundle.name}/${bundle.target}\n"
       "    variables:\n"
       "      catalog: training\n"
       "      schema: staging\n"
       "    permissions:\n"
       "      - user_name: ${workspace.current_user.userName}\n"
       "        level: CAN_MANAGE",
       csize=11, ch=4.72)

s_bullets(nx(), "targets.<t>.resources merges into the base",
          "A target may carry a partial resource definition", [
    ("It is merged, not substituted",
     "The base declaration is read first. The target's fields are then laid over it,\n"
     "field by field, at every level of nesting."),
    ("So a target block names only what differs",
     "Not the job. Not the task list. Not the fields of the cluster that are the same\n"
     "in both. Only the ones that change."),
    ("And the base file is never edited to add a target",
     "resources/ingest.job.yml is written once, in Part D, and stays as it is for the rest\n"
     "of the course."),
], size=13.5, gap=15)

s_runsee(nx(), "The merge, resolved",
         "Staging declares one key. Read what the CLI resolves.",
         "# databricks.yml\n"
         "targets:\n"
         "  staging:\n"
         "    resources:\n"
         "      jobs:\n"
         "        ingest:\n"
         "          job_clusters:\n"
         "            - job_cluster_key: main\n"
         "              new_cluster:\n"
         "                autoscale:\n"
         "                  min_workers: 2\n"
         "                  max_workers: 6\n"
         "\n"
         "$ databricks bundle validate \\\n"
         "    -t staging -o json | jq \\\n"
         "    '.resources.jobs.ingest\n"
         "      .job_clusters[0]'",
         "{\n"
         "  \"job_cluster_key\": \"main\",\n"
         "  \"new_cluster\": {\n"
         "    \"autoscale\": {\n"
         "      \"max_workers\": 6,\n"
         "      \"min_workers\": 2\n"
         "    },\n"
         "    \"data_security_mode\":\n"
         "       \"DATA_SECURITY_MODE_AUTO\",\n"
         "    \"node_type_id\":\n"
         "       \"Standard_D3_v2\",\n"
         "    \"spark_version\":\n"
         "       \"16.4.x-scala2.12\"\n"
         "  }\n"
         "}\n"
         "\n"
         "three fields from the base,\n"
         "one from the target",
         csize=10)

s_runsee(nx(), "Sequences merge by key, and a typo appends",
         "The same override with mian instead of main",
         "# databricks.yml\n"
         "targets:\n"
         "  staging:\n"
         "    resources:\n"
         "      jobs:\n"
         "        ingest:\n"
         "          job_clusters:\n"
         "            - job_cluster_key: mian\n"
         "              new_cluster:\n"
         "                autoscale:\n"
         "                  min_workers: 2\n"
         "                  max_workers: 6\n"
         "\n"
         "$ databricks bundle validate \\\n"
         "    -t staging\n"
         "\n"
         "Validation OK!",
         "job_clusters now has FOUR\n"
         "entries, not three:\n"
         "\n"
         "  main      autoscale 1-3,\n"
         "            untouched, and\n"
         "            still claimed by\n"
         "            land, publish and\n"
         "            notify\n"
         "  features  4 workers\n"
         "  quality   2 workers\n"
         "  mian      autoscale 2-6 and\n"
         "            nothing else: no\n"
         "            runtime, no node\n"
         "            type, claimed by\n"
         "            no task\n"
         "\n"
         "No warning. No error. Staging\n"
         "deploys, runs on 1-3 workers,\n"
         "and creates a cluster that\n"
         "nothing uses.",
         see_ok=False, csize=10)

s_bullets(nx(), "Why that override was silent",
          "The merge matched on a key, and the key it was given did not exist", [
    ("Sequence entries are matched by their key field",
     "job_clusters on job_cluster_key, tasks on task_key. A matching key merges into\n"
     "the entry. A key that matches nothing is a new entry, which is a legitimate\n"
     "thing for a target to add."),
    ("So the CLI cannot tell your typo from your intention",
     "Adding a staging-only task is the same operation as misspelling an existing one.\n"
     "There is no warning available to give."),
    ("The check is to read the resolved configuration, not the file",
     "bundle validate -o json for the target, and count the entries. This is the one\n"
     "fault in the module that no message will ever point at."),
], size=13, gap=12)

s_table(nx(), "Where a variable's value comes from",
        "Five sources, lowest precedence first. Measured, not quoted.",
        ["#", "Source", "Set by"],
        [["1", "default: in the variable declaration", "the author of the bundle"],
         ["2", "the selected target's variables: block", "the author, per environment"],
         ["3", ".databricks/bundle/<target>/variable-overrides.json", "a developer, locally"],
         ["4", "the BUNDLE_VAR_<name> environment variable", "a shell, or a CI runner"],
         ["5", "--var name=value on the command line", "whoever typed the command"]],
        [0.60, 6.15, 3.65], rowh=0.40, size=10.5, mono_cols=(1,),
        note="Correction to Module 3, slide S24: the prefix is BUNDLE_VAR_, not "
             "DATABRICKS_BUNDLE_VAR_, and there are five sources rather than four. --var is "
             "parsed as comma-separated, so a value containing a comma cannot go there at all.")

s_code(nx(), "A variable may hold a whole mapping", "type: complex, and the job references it as one value",
       "variables:\n"
       "  cluster:\n"
       "    description: The job cluster specification for this target\n"
       "    type: complex\n"
       "    default:\n"
       "      spark_version: 16.4.x-scala2.12\n"
       "      node_type_id: Standard_D3_v2\n"
       "      data_security_mode: DATA_SECURITY_MODE_AUTO\n"
       "      autoscale: { min_workers: 1, max_workers: 4 }\n"
       "\n"
       "# resources/ingest.job.yml\n"
       "      job_clusters:\n"
       "        - job_cluster_key: main\n"
       "          new_cluster: ${var.cluster}\n"
       "\n"
       "# databricks.yml, in the staging target\n"
       "    variables:\n"
       "      cluster:\n"
       "        spark_version: 16.4.x-scala2.12\n"
       "        node_type_id: Standard_D3_v2\n"
       "        num_workers: 8",
       csize=10.5, ch=4.96)

s_table(nx(), "Merge or replace",
        "The two mechanisms behave differently, so the choice is deliberate",
        ["Mechanism", "A partial value", "Choose it when"],
        [["targets.<t>.resources",
          "merges field by field\ninto the base",
          "one or two fields differ, and the\nrest must stay in one place"],
         ["a complex variable",
          "replaces the whole value",
          "the specification differs as a unit,\nand a partial copy would mislead"]],
        [3.05, 3.25, 4.10], rowh=0.32, size=10.5, mono_cols=(0,),
        note="The staging assignment on the previous slide carries num_workers and no autoscale, "
             "and the resolved cluster has num_workers and no autoscale. Nothing was inherited. "
             "That is the whole difference between the two rows.")

s_challenge(nx(), 4, "Promote to staging", "10 MIN", [
    "Add a staging target that differs from dev in the catalog, the worker count and "
    "whether the schedule fires. Do not edit resources/ingest.job.yml.",
    "Resolve both targets to JSON and diff them. Account for every line that differs, "
    "including the ones you did not write.",
    "Count the entries in job_clusters for staging. Say why counting is necessary here "
    "and nowhere else.",
    "For each of the two variables, name which of the five sources supplied its value "
    "in staging, and which supplied it in dev.",
], "Two resolutions differing in exactly the properties you declared, and no extra cluster.")

# --- Part F ------------------------------------------------------------------
s_part(nx(), "PART F", "The loop, and the diagnosis",
       "Five steps, each able to fail for reasons the previous one cannot see. Then "
       "the recipe that turns a message into the file you have to open.")

s_table(nx(), "The five steps, and what each one can tell you",
        "Each reads the output of the one before it",
        ["Step", "Reads", "Can report"],
        [["edit", "nothing", "your editor's YAML errors, and no more"],
         ["validate", "the files, and the workspace", "schema faults, unresolved references"],
         ["plan", "the above, and deployment state", "what would be created, changed, deleted"],
         ["deploy", "the above, and the Jobs API", "build, upload, resource and state faults"],
         ["run", "the deployed job", "per-task states, and your code's failures"]],
        [1.60, 3.30, 5.50], rowh=0.42, size=10.5, mono_cols=(0,),
        note="The step that reported a failure bounds the set of possible causes before you read "
             "anything else. That is the first question in the recipe two slides from here.")

s_shot(nx(), "plan before deploy", "The action the deployment would take, per resource",
       "S54", "$ databricks bundle plan -t dev",
       "A create for each resource on the first run; no action for a resource that has not "
       "changed. Module 7 runs this same command inside a pull request, where the plan is "
       "what a reviewer reads.")

s_shot(nx(), "Deploying a second time", "The same command, with nothing changed",
       "S55", "$ databricks bundle deploy -t dev\n"
              "$ databricks bundle deploy -t dev      # again, no edits",
       "Idempotence is a property of the deployment state under state/, not of the "
       "declaration. Delete the state and the second deployment is a first deployment.")

s_shot(nx(), "bundle run, and what it reports", "The run URL first, then the states, then the result",
       "S56", "$ databricks bundle run ingest -t dev\n"
              "$ databricks bundle run ingest -t dev --only features",
       "--only reruns one task of a deployed job without redeploying. Arguments after -- are "
       "passed through to the task, which is how you vary a parameter without editing the file.")

s_challenge(nx(), 5, "Work the loop", "8 MIN", [
    "Break the file that your second task runs: an unclosed parenthesis is enough.",
    "Predict, on paper, the state each of your tasks will be reported in. Then deploy "
    "and run.",
    "Read the run history and record the actual state of every task. Note which states "
    "your code produced and which the platform produced.",
    "Fix the file, deploy again, and read what plan reports for the resources that did "
    "not change.",
], "A failed task, a downstream task reported separately rather than as a failure, and a "
   "second plan that reports no action where nothing changed.")

s_bullets(nx(), "The diagnosis recipe",
          "Four questions, in this order. Each eliminates one class of cause.", [
    ("1. Which command failed: validate, plan, deploy or run?",
     "This bounds the possible causes before you read the message. A failure at run is\n"
     "not a configuration fault, and a failure at validate is not a workspace fault."),
    ("2. If deploy: which stage does the message name?",
     "Build, upload, resource, or state. The stage names the cause: build is your\n"
     "toolchain, upload is permissions or paths, resource is the Jobs API rejecting\n"
     "the definition, state is a conflict with a previous deployment."),
    ("3. Does the message carry a configuration path and a file position?",
     "If it does, the fault is in the declaration and the position is the answer. Stop\n"
     "reading and go there."),
    ("4. If it came from the workspace rather than the CLI, what did we send it?",
     "Read the field out of bundle validate -o json, not out of the file. The file is\n"
     "not what was sent."),
], size=12.5, gap=8)

s_table(nx(), "Failure modes and where each is detected",
        "Ordered by the stage that reports them",
        ["Symptom", "Cause", "Detected at"],
        [["job_cluster_key <k> is not defined", "a task claims a cluster nothing declares", "validate · warning"],
         ["unknown field: <key>", "a misspelled key", "validate · warning"],
         ["invalid value for enum field", "a value outside run_if's six", "validate · warning"],
         ["notebook <path> not found", "the path names no file", "validate · error"],
         ["expected a notebook but got a file", "the marker line is missing", "validate · error"],
         ["no value assigned to required variable", "declared, never assigned, no default", "validate · error"],
         ["a cluster appears that no task claims", "a typo in a target override appended", "*nothing — count them"],
         ["deploy fails after Validation OK!", "an invalid graph, or conflicting compute keys", "deploy · resource"],
         ["a task fails, and those below do not run", "the declaration is right; the code is not", "run"]],
        [4.30, 4.10, 2.00], rowh=0.30, size=10, mono_cols=(0,))

s_challenge(nx(), 6, "Timed team challenge — a sabotaged bundle", "12 MIN", [
    "In teams of three. Copy assets/broken-bundle/ into a scratch directory. It carries "
    "four faults, one from each detection stage on the previous slide.",
    "Work the stages in order. Do not read the whole file first; let the tooling tell "
    "you where to look.",
    "Record the order in which you found them, and which command found each.",
    "Repair all four, and get the run to a terminal success state.",
    "Hands up when the run succeeds. The debrief is the order, not the fixes.",
], "validate clean, deploy successful, run successful, and a written order of discovery.")

# --- Part G ------------------------------------------------------------------
s_bullets(nx(), "What this module did not cover",
          "Four things that belong to later modules, and why", [
    ("Tests over the code these tasks run — Module 5",
     "Everything today verified the declaration. Nothing verified the transformation\n"
     "inside features.py."),
    ("Running validate and deploy on a server — Module 7",
     "Every command today authenticated as you. A runner cannot, which is what the\n"
     "service principal is for."),
    ("An approval gate between staging and production — Module 7",
     "Part E promoted a bundle by typing a different -t. A promotion nobody approved is\n"
     "not a promotion."),
    ("Integration and smoke tests against a deployed bundle — Module 8",
     "Challenge 5 confirmed that a task failed. It did not confirm that a task that\n"
     "succeeded did the right thing."),
], size=12.5, gap=10)

s_bullets(nx(), "Exit ticket", "Four questions. Two minutes. On paper.", [
    (None, "1.  Name the two keys a job_clusters entry requires, and the key on a task that\n"
           "     claims one of them."),
    (None, "2.  A task fails. Name the state the run reports for the task below it, and say\n"
           "     which of the two states your code produced."),
    (None, "3.  Name the three sources of a variable's value that outrank the assignment in\n"
           "     the selected target."),
    (None, "4.  A target override misspells a job_cluster_key. State what appears in the\n"
           "     resolved configuration, and what command reports it."),
], size=13, gap=20)

s_end(nx())

TOTAL = i

# substitute the page-count token now that the total is known
for sl in prs.slides:
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "%%TOTAL%%" in r.text:
                    r.text = r.text.replace("%%TOTAL%%", str(TOTAL))

prs.save("slides.pptx")
print(f"{TOTAL} slides written to slides.pptx")
